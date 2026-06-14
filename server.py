import asyncio
import ipaddress
import json
import os
import re
import shutil
import socket
import sqlite3
import subprocess
import threading
import time
import urllib.error
import urllib.request
import uuid
from collections import defaultdict
from contextlib import asynccontextmanager, contextmanager
from html.parser import HTMLParser
from pathlib import Path
from typing import Dict, Iterator, List, Optional, Set
from urllib.parse import urljoin, urlparse

from fastapi import FastAPI, File, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, Response, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from claude_web import __version__

_PKG_DIR = Path(__file__).parent
_DATA_DIR = Path(os.environ.get("CLAUDE_WEB_DATA_DIR", "")).resolve() if os.environ.get("CLAUDE_WEB_DATA_DIR") else Path.cwd()

STATIC_DIR = _PKG_DIR / "static"
HISTORY_DIR = _DATA_DIR / "history"
UPLOADS_DIR = _DATA_DIR / "uploads"
DB_PATH = _DATA_DIR / "claude-web.db"

HISTORY_DIR.mkdir(parents=True, exist_ok=True)
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
MAX_UPLOAD_MB = 20
IGNORED_DIRS = {".git", "node_modules", ".venv", "venv", "__pycache__", ".next", "dist", "build", ".cache", ".idea", ".vscode"}
KNOWN_TOOL_NAMES = {
    "Bash", "Read", "Write", "Edit", "MultiEdit", "Grep", "Glob",
    "WebFetch", "WebSearch", "TodoWrite", "Task", "NotebookEdit",
}

_running_processes: Dict[str, asyncio.subprocess.Process] = {}
_stopped_sessions: Set[str] = set()
_compacting_sessions: Set[str] = set()
_event_locks: Dict[str, threading.Lock] = {}
_event_lock_refs: Dict[str, int] = {}
_event_lock_access: Dict[str, float] = {}
_event_locks_guard = threading.Lock()
_MAX_EVENT_LOCKS = 1024
_stats_backfill_lock: Optional[asyncio.Lock] = None
_stats_backfill_done = False


class ClaudeCliResolutionError(RuntimeError):
    pass


def resolve_claude_cli_command() -> Optional[str]:
    candidates = ["claude"]
    if os.name == "nt":
        # npm on Windows may put both a Unix shim named "claude" and a usable
        # batch shim named "claude.cmd" on PATH. Python can pick the Unix shim
        # first, so prefer Windows-native launchers explicitly.
        candidates = ["claude.cmd", "claude.exe", "claude.bat", "claude"]
    for candidate in candidates:
        path = shutil.which(candidate)
        if path:
            return path
    return None


def claude_cli_command() -> str:
    command = resolve_claude_cli_command()
    if command:
        return command
    return "claude.cmd" if os.name == "nt" else "claude"


def _claude_package_bin(package_dir: Path) -> Optional[Path]:
    package_json = package_dir / "package.json"
    if package_json.exists():
        try:
            data = json.loads(package_json.read_text(encoding="utf-8"))
            bin_entry = data.get("bin")
            if isinstance(bin_entry, dict):
                bin_entry = bin_entry.get("claude") or next(iter(bin_entry.values()), None)
            if isinstance(bin_entry, str):
                candidate = (package_dir / bin_entry).resolve()
                if candidate.exists():
                    return candidate
        except Exception:
            pass
    for name in ("cli.js", "cli.mjs"):
        candidate = package_dir / name
        if candidate.exists():
            return candidate.resolve()
    return None


def _windows_claude_node_argv(command: str) -> Optional[List[str]]:
    command_path = Path(command)
    bin_dir = command_path.parent
    package_dirs = [
        bin_dir / "node_modules" / "@anthropic-ai" / "claude-code",
        bin_dir.parent / "@anthropic-ai" / "claude-code",
    ]
    script = next((p for p in (_claude_package_bin(d) for d in package_dirs) if p), None)
    if script is None:
        return None

    # claude-code 2.x ships a native Windows launcher (bin/claude.exe); invoke
    # it directly. node.exe can't load an .exe as a JS module.
    if script.suffix.lower() in (".exe", ".com"):
        return [str(script)]

    node_candidates = [
        bin_dir / "node.exe",
        shutil.which("node.exe"),
        shutil.which("node"),
    ]
    node = next((str(p) for p in node_candidates if p and Path(p).exists()), None)
    if node is None:
        return None
    return [node, str(script)]


def claude_cli_argv(*args: str, allow_batch_shim: bool = False) -> List[str]:
    command = resolve_claude_cli_command()
    if command is None:
        return ["claude.cmd" if os.name == "nt" else "claude", *args]
    if os.name == "nt" and command.lower().endswith((".cmd", ".bat")):
        node_argv = _windows_claude_node_argv(command)
        if node_argv:
            return [*node_argv, *args]
        if not allow_batch_shim:
            raise ClaudeCliResolutionError(
                "claude CLI batch shim found, but the Node.js entrypoint could not be resolved"
            )
    return [command, *args]


async def _terminate_process(process: asyncio.subprocess.Process, grace: float = 3.0) -> None:
    if process.returncode is not None:
        return
    try:
        process.terminate()
    except ProcessLookupError:
        return
    try:
        await asyncio.wait_for(process.wait(), timeout=grace)
        return
    except asyncio.TimeoutError:
        pass
    try:
        process.kill()
    except ProcessLookupError:
        return
    try:
        await asyncio.wait_for(process.wait(), timeout=2.0)
    except asyncio.TimeoutError:
        pass


async def _shutdown_terminate_running_processes() -> None:
    if not _running_processes:
        return
    processes = list(_running_processes.values())
    _running_processes.clear()
    await asyncio.gather(
        *(_terminate_process(p) for p in processes),
        return_exceptions=True,
    )


_UPLOAD_RETENTION_SECONDS = 30 * 24 * 60 * 60  # 30 days


def _prune_old_uploads(retention_seconds: int = _UPLOAD_RETENTION_SECONDS) -> int:
    """Delete files in UPLOADS_DIR older than retention_seconds. Returns count
    of files removed. Best-effort: silently skips entries we can't stat/unlink."""
    if not UPLOADS_DIR.exists():
        return 0
    cutoff = time.time() - retention_seconds
    removed = 0
    for entry in UPLOADS_DIR.iterdir():
        if not entry.is_file():
            continue
        try:
            if entry.stat().st_mtime < cutoff:
                entry.unlink()
                removed += 1
        except OSError:
            continue
    return removed


@asynccontextmanager
async def _lifespan(app: FastAPI):
    # Prune stale uploads in a background thread so startup isn't blocked on disk IO.
    asyncio.get_event_loop().run_in_executor(None, _prune_old_uploads)
    try:
        yield
    finally:
        await _shutdown_terminate_running_processes()


app = FastAPI(title="Claude Code Web", lifespan=_lifespan)


async def _drain_stream(stream: asyncio.StreamReader, buffer: bytearray, limit: int = 256 * 1024) -> None:
    try:
        while True:
            chunk = await stream.read(8192)
            if not chunk:
                return
            remaining = limit - len(buffer)
            if remaining > 0:
                buffer.extend(chunk[:remaining])
    except asyncio.CancelledError:
        raise
    except Exception:
        return


_DB_INITIALIZED = False


@contextmanager
def db_connect() -> Iterator[sqlite3.Connection]:
    global _DB_INITIALIZED
    conn = sqlite3.connect(DB_PATH, timeout=10)
    try:
        conn.row_factory = sqlite3.Row
        if not _DB_INITIALIZED:
            conn.execute("PRAGMA journal_mode=WAL")
            conn.execute("PRAGMA synchronous=NORMAL")
            _DB_INITIALIZED = True
        conn.execute("PRAGMA busy_timeout=5000")
        yield conn
        conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        conn.close()


def ensure_column(conn: sqlite3.Connection, table: str, column: str, ddl: str) -> None:
    cols = {row["name"] for row in conn.execute(f"PRAGMA table_info({table})").fetchall()}
    if column not in cols:
        conn.execute(f"ALTER TABLE {table} ADD COLUMN {column} {ddl}")


@contextmanager
def session_event_lock(session_id: str) -> Iterator[None]:
    now = time.time()
    with _event_locks_guard:
        lock = _event_locks.get(session_id)
        if lock is None:
            lock = threading.Lock()
            _event_locks[session_id] = lock
            _event_lock_refs[session_id] = 0
        _event_lock_refs[session_id] = _event_lock_refs.get(session_id, 0) + 1
        _event_lock_access[session_id] = now
    lock.acquire()
    try:
        yield
    finally:
        lock.release()
        with _event_locks_guard:
            _event_lock_refs[session_id] = max(_event_lock_refs.get(session_id, 1) - 1, 0)
            _event_lock_access[session_id] = time.time()
            prune_event_locks_locked()


def prune_event_locks_locked() -> None:
    if len(_event_locks) <= _MAX_EVENT_LOCKS:
        return
    removable = [
        (last_access, sid)
        for sid, last_access in _event_lock_access.items()
        if _event_lock_refs.get(sid, 0) == 0
    ]
    removable.sort()
    for _, sid in removable[: max(1, len(_event_locks) - _MAX_EVENT_LOCKS)]:
        _event_locks.pop(sid, None)
        _event_lock_refs.pop(sid, None)
        _event_lock_access.pop(sid, None)


def init_db() -> None:
    with db_connect() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS sessions (
                id TEXT PRIMARY KEY,
                title TEXT NOT NULL DEFAULT '',
                cwd TEXT NOT NULL DEFAULT '',
                created_at REAL NOT NULL,
                updated_at REAL NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS prompts (
                id TEXT PRIMARY KEY,
                name TEXT NOT NULL,
                content TEXT NOT NULL,
                created_at REAL NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS session_usage (
                id TEXT PRIMARY KEY,
                session_id TEXT NOT NULL,
                turn_idx INTEGER NOT NULL,
                input_tokens INTEGER NOT NULL DEFAULT 0,
                output_tokens INTEGER NOT NULL DEFAULT 0,
                cache_read_input_tokens INTEGER NOT NULL DEFAULT 0,
                cache_creation_input_tokens INTEGER NOT NULL DEFAULT 0,
                total_cost_usd REAL NOT NULL DEFAULT 0,
                ts REAL NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS memories (
                id TEXT PRIMARY KEY,
                content TEXT NOT NULL,
                enabled INTEGER NOT NULL DEFAULT 1,
                scope TEXT NOT NULL DEFAULT 'global',
                created_at REAL NOT NULL,
                updated_at REAL NOT NULL
            )
            """
        )
        ensure_column(conn, "sessions", "pinned", "INTEGER NOT NULL DEFAULT 0")
        ensure_column(conn, "sessions", "archived", "INTEGER NOT NULL DEFAULT 0")
        ensure_column(conn, "sessions", "tags", "TEXT NOT NULL DEFAULT ''")
        ensure_column(conn, "sessions", "manual_title", "INTEGER NOT NULL DEFAULT 0")
        ensure_column(conn, "sessions", "remote_session_id", "TEXT NOT NULL DEFAULT ''")
        ensure_column(conn, "sessions", "remote_ready", "INTEGER NOT NULL DEFAULT 0")
        ensure_column(conn, "sessions", "summary_cache", "TEXT")
        ensure_column(conn, "prompts", "slash_trigger", "TEXT NOT NULL DEFAULT ''")
        ensure_column(conn, "session_usage", "duration_ms", "REAL NOT NULL DEFAULT 0")
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS tool_calls (
                id TEXT PRIMARY KEY,
                session_id TEXT NOT NULL,
                tool_name TEXT NOT NULL,
                ts REAL NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS app_meta (
                key TEXT PRIMARY KEY,
                value TEXT NOT NULL
            )
            """
        )
        conn.execute("CREATE INDEX IF NOT EXISTS idx_session_usage_session ON session_usage(session_id, turn_idx)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_session_usage_ts ON session_usage(ts)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_tool_calls_session ON tool_calls(session_id)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_tool_calls_name ON tool_calls(tool_name)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_memories_scope ON memories(scope, enabled)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_sessions_summary_cache ON sessions(summary_cache)")
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS cwd_history (
                id TEXT PRIMARY KEY,
                path TEXT NOT NULL UNIQUE,
                created_at REAL NOT NULL
            )
            """
        )


init_db()


def upsert_session(session_id: str, title: str, cwd: str) -> None:
    now = time.time()
    with db_connect() as conn:
        row = conn.execute(
            "SELECT title, manual_title FROM sessions WHERE id = ?", (session_id,)
        ).fetchone()
        if row is None:
            conn.execute(
                "INSERT INTO sessions (id, title, cwd, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
                (session_id, title, cwd, now, now),
            )
        else:
            new_title = row["title"]
            if not row["manual_title"] and not new_title:
                new_title = title
            conn.execute(
                "UPDATE sessions SET title = ?, cwd = ?, updated_at = ? WHERE id = ?",
                (new_title, cwd, now, session_id),
            )


_SUMMARY_CACHE_LIMIT = 20000


def trim_summary_cache(text: str) -> str:
    return text[-_SUMMARY_CACHE_LIMIT:]


def summarize_cache_from_events(events: List[dict]) -> str:
    return trim_summary_cache(summarize_text_from_events(events))


def set_session_summary_cache(session_id: str, summary: str) -> None:
    with db_connect() as conn:
        conn.execute("UPDATE sessions SET summary_cache = ? WHERE id = ?", (summary, session_id))


def update_session_summary_cache_for_event(conn: sqlite3.Connection, session_id: str, event: dict) -> None:
    snippet = summarize_text_from_events([event]).strip()
    if not snippet:
        return
    conn.execute(
        """
        UPDATE sessions
        SET summary_cache = substr(COALESCE(summary_cache, '') || ? || char(10), ?)
        WHERE id = ?
        """,
        (snippet, -_SUMMARY_CACHE_LIMIT, session_id),
    )


def ensure_session_summary_cache(session_id: str, current_summary: Optional[str]) -> str:
    if current_summary is not None:
        return current_summary
    events = load_events(session_id)
    summary = summarize_cache_from_events(events)
    set_session_summary_cache(session_id, summary)
    return summary


def tool_call_rows_from_event(session_id: str, event: dict) -> List[tuple]:
    if event.get("type") != "assistant":
        return []
    content = (event.get("message") or {}).get("content") or []
    names = [
        block.get("name") or "?"
        for block in content
        if isinstance(block, dict) and block.get("type") == "tool_use"
    ]
    if not names:
        return []
    now = float(event.get("ts") or time.time())
    return [(uuid.uuid4().hex, session_id, name, now) for name in names]


def insert_tool_call_rows(conn: sqlite3.Connection, rows: List[tuple]) -> None:
    if not rows:
        return
    conn.executemany(
        "INSERT INTO tool_calls (id, session_id, tool_name, ts) VALUES (?, ?, ?, ?)",
        rows,
    )


def replace_session_tool_call_rows(conn: sqlite3.Connection, session_id: str, events: List[dict]) -> None:
    conn.execute("DELETE FROM tool_calls WHERE session_id = ?", (session_id,))
    rows: List[tuple] = []
    for event in events:
        rows.extend(tool_call_rows_from_event(session_id, event))
    insert_tool_call_rows(conn, rows)


def record_tool_calls(session_id: str, event: dict) -> None:
    rows = tool_call_rows_from_event(session_id, event)
    if not rows:
        return
    with db_connect() as conn:
        insert_tool_call_rows(conn, rows)


def append_event(session_id: str, event: dict) -> None:
    path = HISTORY_DIR / f"{session_id}.jsonl"
    with session_event_lock(session_id):
        with path.open("a", encoding="utf-8") as f:
            f.write(json.dumps(event, ensure_ascii=False) + "\n")
        with db_connect() as conn:
            update_session_summary_cache_for_event(conn, session_id, event)
            insert_tool_call_rows(conn, tool_call_rows_from_event(session_id, event))


def record_usage(session_id: str, result_event: dict) -> None:
    usage = result_event.get("usage") or {}
    if not isinstance(usage, dict):
        return
    input_tokens = int(usage.get("input_tokens") or 0)
    output_tokens = int(usage.get("output_tokens") or 0)
    cache_read = int(usage.get("cache_read_input_tokens") or 0)
    cache_create = int(usage.get("cache_creation_input_tokens") or 0)
    cost = float(result_event.get("total_cost_usd") or 0)
    duration_ms = float(result_event.get("duration_ms") or 0)
    if input_tokens == 0 and output_tokens == 0 and cache_read == 0 and cache_create == 0 and cost == 0:
        return
    with db_connect() as conn:
        conn.execute(
            """
            INSERT INTO session_usage (
                id, session_id, turn_idx, input_tokens, output_tokens,
                cache_read_input_tokens, cache_creation_input_tokens,
                total_cost_usd, duration_ms, ts
            ) VALUES (
                ?, ?,
                COALESCE((SELECT MAX(turn_idx) FROM session_usage WHERE session_id = ?), 0) + 1,
                ?, ?, ?, ?, ?, ?, ?
            )
            """,
            (
                uuid.uuid4().hex, session_id, session_id, input_tokens, output_tokens,
                cache_read, cache_create, cost, duration_ms, time.time(),
            ),
        )


def normalize_memory_scope(scope: Optional[str]) -> str:
    raw_scope = (scope or "global").strip() or "global"
    if raw_scope.startswith("project:"):
        raw_path = raw_scope[len("project:") :].strip()
        if raw_path:
            return "project:" + str(Path(os.path.expanduser(raw_path)).resolve())
        return "global"
    if raw_scope.startswith("session:") and raw_scope[len("session:") :].strip():
        return raw_scope
    if raw_scope == "global":
        return "global"
    return raw_scope


def matching_memory_scopes(cwd: str, session_id: str) -> List[str]:
    scopes = ["global"]
    if cwd:
        scopes.append(normalize_memory_scope(f"project:{cwd}"))
    if session_id:
        scopes.append(f"session:{session_id}")
    return scopes


def load_enabled_memories(cwd: str, session_id: str) -> List[dict]:
    scopes = matching_memory_scopes(cwd, session_id)
    placeholders = ",".join("?" for _ in scopes)
    with db_connect() as conn:
        rows = conn.execute(
            f"""
            SELECT id, content, enabled, scope, created_at, updated_at
            FROM memories
            WHERE enabled = 1 AND scope IN ({placeholders})
            ORDER BY scope, updated_at DESC
            """,
            scopes,
        ).fetchall()
    return [dict(r) for r in rows]


def compose_system_prompt(memory_items: List[dict], user_system_prompt: Optional[str]) -> Optional[str]:
    parts: List[str] = []
    if memory_items:
        memory_text = "\n".join(f"- {m['content']}" for m in memory_items if m.get("content"))
        parts.append("Persistent memory for this user/project/session:\n" + memory_text)
    if user_system_prompt:
        parts.append(user_system_prompt)
    return "\n\n".join(parts) if parts else None


def load_events(session_id: str) -> List[dict]:
    path = HISTORY_DIR / f"{session_id}.jsonl"
    if not path.exists():
        return []
    events: List[dict] = []
    with path.open("r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                events.append(json.loads(line))
            except json.JSONDecodeError:
                continue
    return events


def iter_history_paths() -> Iterator[Path]:
    for path in HISTORY_DIR.glob("*.jsonl"):
        if ".before-compact-" in path.name or ".tmp." in path.name:
            continue
        yield path


def backfill_tool_calls_once() -> None:
    try:
        with db_connect() as conn:
            row = conn.execute("SELECT value FROM app_meta WHERE key = 'tool_calls_backfilled_v1'").fetchone()
            if row is not None:
                return
            conn.execute("DELETE FROM tool_calls")
            rows: List[tuple] = []
            for path in iter_history_paths():
                session_id = path.stem
                for event in load_events(session_id):
                    rows.extend(tool_call_rows_from_event(session_id, event))
                    if len(rows) >= 1000:
                        insert_tool_call_rows(conn, rows)
                        rows = []
            insert_tool_call_rows(conn, rows)
            conn.execute(
                "INSERT OR REPLACE INTO app_meta (key, value) VALUES ('tool_calls_backfilled_v1', ?)",
                (str(time.time()),),
            )
    except Exception:
        return


def backfill_usage_duration_once() -> None:
    try:
        with db_connect() as conn:
            row = conn.execute("SELECT value FROM app_meta WHERE key = 'usage_duration_backfilled_v1'").fetchone()
            if row is not None:
                return
            for path in iter_history_paths():
                session_id = path.stem
                result_events = [event for event in load_events(session_id) if event.get("type") == "result"]
                if not result_events:
                    continue
                usage_rows = conn.execute(
                    """
                    SELECT id, duration_ms
                    FROM session_usage
                    WHERE session_id = ?
                    ORDER BY turn_idx
                    """,
                    (session_id,),
                ).fetchall()
                for usage_row, event in zip(usage_rows, result_events):
                    duration_ms = float(event.get("duration_ms") or 0)
                    if duration_ms > 0 and float(usage_row["duration_ms"] or 0) == 0:
                        conn.execute(
                            "UPDATE session_usage SET duration_ms = ? WHERE id = ?",
                            (duration_ms, usage_row["id"]),
                        )
            conn.execute(
                "INSERT OR REPLACE INTO app_meta (key, value) VALUES ('usage_duration_backfilled_v1', ?)",
                (str(time.time()),),
            )
    except Exception:
        return


async def ensure_stats_backfilled() -> None:
    global _stats_backfill_done, _stats_backfill_lock
    if _stats_backfill_done:
        return
    if _stats_backfill_lock is None:
        _stats_backfill_lock = asyncio.Lock()
    if _stats_backfill_lock.locked():
        return
    async with _stats_backfill_lock:
        if _stats_backfill_done:
            return
        await asyncio.to_thread(backfill_usage_duration_once)
        await asyncio.to_thread(backfill_tool_calls_once)
        _stats_backfill_done = True


def save_events(session_id: str, events: List[dict]) -> None:
    path = HISTORY_DIR / f"{session_id}.jsonl"
    with session_event_lock(session_id):
        if not events:
            if path.exists():
                path.unlink()
            with db_connect() as conn:
                conn.execute("UPDATE sessions SET summary_cache = ? WHERE id = ?", ("", session_id))
                conn.execute("DELETE FROM tool_calls WHERE session_id = ?", (session_id,))
            return
        tmp_path = path.with_suffix(path.suffix + f".tmp.{os.getpid()}.{uuid.uuid4().hex[:8]}")
        try:
            with tmp_path.open("w", encoding="utf-8") as f:
                for event in events:
                    f.write(json.dumps(event, ensure_ascii=False) + "\n")
                f.flush()
                os.fsync(f.fileno())
            os.replace(tmp_path, path)
        except Exception:
            tmp_path.unlink(missing_ok=True)
            raise
        with db_connect() as conn:
            conn.execute("UPDATE sessions SET summary_cache = ? WHERE id = ?", (summarize_cache_from_events(events), session_id))
            replace_session_tool_call_rows(conn, session_id, events)


def summarize_text_from_events(events: List[dict]) -> str:
    parts: List[str] = []
    for ev in events:
        if ev.get("type") == "user_input":
            parts.append(ev.get("text", ""))
        elif ev.get("type") == "assistant":
            content = (ev.get("message") or {}).get("content") or []
            for block in content:
                if block.get("type") == "text":
                    parts.append(block.get("text", ""))
    return "\n".join(parts)


class ChatRequest(BaseModel):
    message: str
    session_id: Optional[str] = None
    cwd: Optional[str] = None
    images: Optional[List[str]] = None
    model: Optional[str] = None
    system_prompt: Optional[str] = None
    display_message: Optional[str] = None
    permission_mode: Optional[str] = None
    allowed_tools: Optional[List[str]] = None
    disallowed_tools: Optional[List[str]] = None
    force_new: Optional[bool] = None
    # UI-only metadata for attached docs (name/size/length/path); rendered as
    # badges on the user message. Not used to build the prompt — the doc text
    # is already embedded in `message` by the client.
    docs: Optional[List[dict]] = None


class PromptRequest(BaseModel):
    name: str
    content: str
    slash_trigger: Optional[str] = ""


class MemoryRequest(BaseModel):
    content: str
    enabled: Optional[bool] = True
    scope: Optional[str] = "global"


class SessionPatch(BaseModel):
    title: Optional[str] = None
    pinned: Optional[bool] = None
    archived: Optional[bool] = None
    tags: Optional[str] = None


class ForkRequest(BaseModel):
    event_index: int
    new_text: Optional[str] = None


class RestoreRequest(BaseModel):
    event_index: int


class FetchUrlRequest(BaseModel):
    url: str
    max_chars: Optional[int] = 10000


def build_args(
    message: str,
    session_id: str,
    resume: bool,
    model: Optional[str],
    system_prompt: Optional[str],
    permission_mode: Optional[str] = None,
    allowed_tools: Optional[List[str]] = None,
    disallowed_tools: Optional[List[str]] = None,
    use_stdin: bool = False,
) -> List[str]:
    args = claude_cli_argv()
    if use_stdin:
        args += ["-p", "--input-format", "stream-json"]
    else:
        args += ["-p", message]
    args += [
        "--output-format", "stream-json",
        "--verbose",
        "--include-partial-messages",
    ]
    if resume:
        args += ["--resume", session_id]
    else:
        args += ["--session-id", session_id]
    if model:
        args += ["--model", model]
    if system_prompt:
        args += ["--append-system-prompt", system_prompt]
    if permission_mode and permission_mode in ("default", "acceptEdits", "bypassPermissions", "plan"):
        args += ["--permission-mode", permission_mode]
    if allowed_tools:
        args += ["--allowed-tools", ",".join(allowed_tools)]
    if disallowed_tools:
        args += ["--disallowed-tools", ",".join(disallowed_tools)]
    return args


def extract_tool_name(text: str) -> Optional[str]:
    mcp_match = re.search(r"\bmcp__[A-Za-z0-9_:-]+(?:__[A-Za-z0-9_:-]+)*\b", text)
    if mcp_match:
        return mcp_match.group(0)

    for tool in KNOWN_TOOL_NAMES:
        if re.search(rf"\b{re.escape(tool)}\b", text):
            return tool

    patterns = [
        r"(?:MCP tool|mcp tool|tool)\s+[\"'`]?([A-Za-z][A-Za-z0-9_:-]{1,80})[\"'`]?",
        r"[\"'`]([A-Za-z][A-Za-z0-9_:-]{1,80})[\"'`]\s+(?:tool|Tool|MCP tool|mcp tool)",
    ]
    stop_words = {"approval", "permission", "tool", "tools", "mcp", "required", "requires", "non-interactive"}
    for pattern in patterns:
        m = re.search(pattern, text)
        if not m:
            continue
        candidate = m.group(1).strip()
        if candidate.lower() not in stop_words:
            return candidate
    return None


def classify_claude_error(message: str) -> dict:
    text = (message or "").strip() or "claude exited with error"
    lower = text.lower()
    tool_name = extract_tool_name(text)
    permissionish = any(k in lower for k in (
        "requires approval", "approval required", "needs approval", "approval",
        "cannot prompt", "non-interactive", "not allowed", "permission denied",
    ))
    if ("permission" in lower and ("tool" in lower or "mcp" in lower or tool_name)) or (permissionish and ("tool" in lower or "mcp" in lower or tool_name)):
        return {
            "type": "permission_error",
            "message": text,
            "tool_name": tool_name,
            "hint": "当前 Web UI 不支持运行中批准工具权限；请预先放行工具后重试本轮，或改用 Claude Code CLI。",
        }
    return {"type": "error", "message": text}


def build_image_input_message(message: str, images: List[str]) -> bytes:
    """Build a stream-json user message. Works with or without images."""
    import base64 as b64mod
    content: List[dict] = []
    for img_path in images:
        p = Path(img_path)
        if not p.exists():
            continue
        ext = p.suffix.lower()
        media_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                     ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp"}
        media_type = media_map.get(ext, "image/png")
        data = b64mod.b64encode(p.read_bytes()).decode()
        content.append({"type": "image", "source": {"type": "base64", "media_type": media_type, "data": data}})
    content.append({"type": "text", "text": message})
    msg = {"type": "user", "message": {"role": "user", "content": content}}
    return json.dumps(msg, ensure_ascii=False).encode() + b"\n"


async def _git_run(cwd: str, *args: str) -> Optional[str]:
    try:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", cwd, *args,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=10)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return None
        if proc.returncode != 0:
            return None
        return stdout.decode("utf-8", errors="replace").strip()
    except Exception:
        return None


async def create_git_checkpoint(cwd: str) -> Optional[dict]:
    if not cwd or not os.path.isdir(cwd):
        return None
    git_dir = await _git_run(cwd, "rev-parse", "--git-dir")
    if git_dir is None:
        return None
    head = await _git_run(cwd, "rev-parse", "HEAD")
    if head is None:
        return None
    stash = await _git_run(cwd, "stash", "create", f"claude-web-checkpoint-{int(time.time())}")
    return {"type": "git", "head": head, "stash": stash or ""}


async def restore_git_checkpoint(cwd: str, cp: dict) -> bool:
    if not cp or cp.get("type") != "git" or not cwd:
        return False
    head = cp.get("head")
    stash = cp.get("stash") or ""
    if not head:
        return False
    if await _git_run(cwd, "reset", "--hard", head) is None:
        return False
    await _git_run(cwd, "clean", "-fd")
    if stash:
        await _git_run(cwd, "stash", "apply", stash)
    return True


def format_context_snippet(events: List[dict], max_chars: int = 6000) -> str:
    lines: List[str] = []
    total = 0
    for ev in events:
        t = ev.get("type")
        if t == "user_input":
            text = (ev.get("text") or "").strip()
            if text:
                chunk = f"用户: {text}"
                lines.append(chunk)
                total += len(chunk)
        elif t == "assistant":
            content = (ev.get("message") or {}).get("content") or []
            for block in content:
                if block.get("type") == "text":
                    text = (block.get("text") or "").strip()
                    if text:
                        chunk = f"助手: {text[:600]}"
                        lines.append(chunk)
                        total += len(chunk)
                elif block.get("type") == "tool_use":
                    name = block.get("name", "")
                    chunk = f"(助手调用了工具: {name})"
                    lines.append(chunk)
                    total += len(chunk)
        if total > max_chars:
            lines.append("...（历史已截断）")
            break
    return "\n\n".join(lines)


def derive_title(message: str) -> str:
    """First line of the user message, with code fences / markdown headers / bullet
    markers stripped so the title reads naturally even when the message starts with
    a code block or markdown."""
    if not message:
        return "未命名会话"
    lines = message.splitlines()
    in_fence = False
    first_in_fence: Optional[str] = None
    for raw in lines:
        stripped = raw.strip()
        if not stripped:
            continue
        # Toggle on triple-backtick fences and skip their content.
        if stripped.startswith("```") or stripped.startswith("~~~"):
            in_fence = not in_fence
            continue
        if in_fence:
            if first_in_fence is None:
                first_in_fence = stripped
            continue
        # Strip markdown header / quote / list prefixes for nicer titles.
        cleaned = re.sub(r"^[#>\-*+\d.\s]+", "", stripped).strip()
        if cleaned:
            return cleaned[:60]
    if first_in_fence:
        return first_in_fence[:60]
    fallback = message.strip().replace("\n", " ")
    return fallback[:60] if fallback else "未命名会话"


def session_has_remote_conversation(events: List[dict]) -> bool:
    for ev in events:
        event_type = ev.get("type")
        if event_type == "user_input" and ev.get("compacted") is True:
            return True
        if event_type == "assistant":
            return True
        if event_type == "system" and ev.get("subtype") == "init":
            return True
        if event_type == "result" and not ev.get("is_error"):
            return True
    return False


def resolve_remote_session_state(session_id: str, row: Optional[sqlite3.Row], events: List[dict]):
    if row is None:
        return session_id, session_has_remote_conversation(events)
    remote_session_id = (row["remote_session_id"] or "").strip() or session_id
    if (row["remote_session_id"] or "").strip():
        return remote_session_id, bool(row["remote_ready"])
    return remote_session_id, session_has_remote_conversation(events)


def set_session_remote_state(session_id: str, remote_session_id: str, remote_ready: bool) -> None:
    now = time.time()
    with db_connect() as conn:
        conn.execute(
            "UPDATE sessions SET remote_session_id = ?, remote_ready = ?, updated_at = ? WHERE id = ?",
            (remote_session_id, 1 if remote_ready else 0, now, session_id),
        )


def prune_session_compact_backups(session_id: str, keep_latest: int = 3, max_age_seconds: int = 7 * 24 * 60 * 60) -> None:
    backups = sorted(
        iter_session_compact_backups(session_id),
        key=lambda x: x.stat().st_mtime,
        reverse=True,
    )
    cutoff = time.time() - max_age_seconds
    for idx, backup in enumerate(backups):
        try:
            if idx >= keep_latest or backup.stat().st_mtime < cutoff:
                backup.unlink(missing_ok=True)
        except OSError:
            continue


def iter_session_compact_backups(session_id: str) -> List[Path]:
    prefix = f"{session_id}.before-compact-"
    try:
        return [
            path for path in HISTORY_DIR.iterdir()
            if path.is_file() and path.name.startswith(prefix) and path.name.endswith(".jsonl")
        ]
    except OSError:
        return []


@app.post("/api/chat")
async def chat(req: ChatRequest):
    session_id = req.session_id or str(uuid.uuid4())
    if session_id in _compacting_sessions:
        raise HTTPException(status_code=409, detail="session is compacting")
    existing_events = load_events(session_id) if req.session_id else []
    with db_connect() as conn:
        row = conn.execute(
            "SELECT cwd, remote_session_id, remote_ready FROM sessions WHERE id = ?",
            (session_id,),
        ).fetchone()

    remote_session_id, remote_ready = resolve_remote_session_state(session_id, row, existing_events)
    if req.force_new is True:
        stored_remote_id = ((row["remote_session_id"] or "").strip() if row else "")
        stored_remote_ready = bool(row["remote_ready"]) if row else False
        if stored_remote_id and not stored_remote_ready:
            remote_session_id = stored_remote_id
        elif row is not None and remote_ready:
            remote_session_id = str(uuid.uuid4())
        remote_ready = False

    is_new = not remote_ready
    work_dir = req.cwd or (row["cwd"] if row and row["cwd"] else os.path.expanduser("~"))
    full_message = req.message
    display_text = req.display_message if req.display_message is not None else req.message

    checkpoint = await create_git_checkpoint(work_dir)

    user_event = {
        "type": "user_input",
        "text": display_text,
        "images": req.images or [],
        "docs": req.docs or [],
        "ts": time.time(),
        "checkpoint": checkpoint,
    }
    # When the prompt was rewritten on the client (doc content / URL fetch / web-search prefix
    # injected), keep the full sent text so badge previews can recover doc bodies even
    # after the upload file is pruned. Only stored when it actually differs.
    if req.message != display_text:
        user_event["full_text"] = req.message
    upsert_session(session_id, derive_title(display_text), work_dir)
    append_event(session_id, user_event)
    set_session_remote_state(session_id, remote_session_id, remote_ready and not is_new)

    async def generate():
        remote_became_ready = remote_ready and not is_new
        meta = {
            "type": "meta",
            "session_id": session_id,
            "cwd": work_dir,
            "has_checkpoint": checkpoint is not None,
        }
        yield f"data: {json.dumps(meta)}\n\n"

        # If a previous chat for the same session is still running (e.g. duplicate
        # request, network retry, fast double-click), terminate it before spawning
        # a new one. Otherwise the old subprocess would be orphaned, burning tokens
        # and producing stray events.
        existing = _running_processes.pop(session_id, None)
        if existing is not None:
            await _terminate_process(existing)
        _stopped_sessions.discard(session_id)

        has_images = bool(req.images)
        # Route through stdin when the prompt would blow past argv limits
        # (macOS ~256KB, Linux ~128KB total argv). Images already force stdin.
        message_too_large = len(full_message.encode("utf-8")) > ARGV_STDIN_THRESHOLD
        use_stdin = has_images or message_too_large
        effective_system_prompt = compose_system_prompt(
            load_enabled_memories(work_dir, session_id),
            req.system_prompt,
        )
        try:
            args = build_args(
                full_message, remote_session_id,
                resume=not is_new,
                model=req.model,
                system_prompt=effective_system_prompt,
                permission_mode=req.permission_mode,
                allowed_tools=req.allowed_tools,
                disallowed_tools=req.disallowed_tools,
                use_stdin=use_stdin,
            )
        except ClaudeCliResolutionError as e:
            err_event = {"type": "error", "message": str(e)}
            append_event(session_id, err_event)
            yield f"data: {json.dumps(err_event, ensure_ascii=False)}\n\n"
            return
        stdin_data: Optional[bytes] = None
        if use_stdin:
            stdin_data = build_image_input_message(full_message, req.images or [])
        try:
            process = await asyncio.create_subprocess_exec(
                *args,
                stdin=asyncio.subprocess.PIPE if use_stdin else None,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                cwd=work_dir,
                limit=16 * 1024 * 1024,
            )
            if use_stdin and stdin_data and process.stdin:
                try:
                    process.stdin.write(stdin_data)
                    await process.stdin.drain()
                    process.stdin.close()
                    await process.stdin.wait_closed()
                except (BrokenPipeError, ConnectionResetError):
                    # CLI exited early (auth failure, bad args, etc). The stderr
                    # path will surface the real reason; don't tear down SSE here.
                    pass
        except FileNotFoundError:
            err_event = {"type": "error", "message": "claude CLI not found in PATH"}
            append_event(session_id, err_event)
            yield f"data: {json.dumps(err_event)}\n\n"
            return

        _running_processes[session_id] = process
        stderr_buffer = bytearray()
        stderr_task: Optional[asyncio.Task] = None
        if process.stderr is not None:
            stderr_task = asyncio.create_task(_drain_stream(process.stderr, stderr_buffer))

        try:
            assert process.stdout is not None
            while True:
                try:
                    raw = await process.stdout.readline()
                except ValueError as e:
                    err_event = {"type": "error", "message": f"stdout line too large: {e}"}
                    append_event(session_id, err_event)
                    yield f"data: {json.dumps(err_event)}\n\n"
                    break
                if not raw:
                    break
                line = raw.decode("utf-8", errors="replace").strip()
                if not line:
                    continue
                try:
                    obj = json.loads(line)
                except json.JSONDecodeError:
                    obj = {"type": "raw", "text": line}
                t = obj.get("type")
                if session_has_remote_conversation([obj]):
                    remote_became_ready = True
                if t != "stream_event" and not (t == "system" and obj.get("subtype", "").startswith("hook_")):
                    append_event(session_id, obj)
                    if t == "result":
                        record_usage(session_id, obj)
                yield f"data: {json.dumps(obj, ensure_ascii=False)}\n\n"

            rc = await process.wait()
            if stderr_task is not None:
                try:
                    await asyncio.wait_for(asyncio.shield(stderr_task), timeout=1.0)
                except asyncio.TimeoutError:
                    pass
            stopped_by_user = session_id in _stopped_sessions
            if rc != 0 and not stopped_by_user:
                err_text = bytes(stderr_buffer).decode("utf-8", errors="replace")
                err_event = classify_claude_error(err_text or f"claude exited with code {rc}")
                append_event(session_id, err_event)
                yield f"data: {json.dumps(err_event, ensure_ascii=False)}\n\n"
        finally:
            if stderr_task is not None and not stderr_task.done():
                stderr_task.cancel()
                try:
                    await stderr_task
                except (asyncio.CancelledError, Exception):
                    pass
            await _terminate_process(process)
            _running_processes.pop(session_id, None)
            _stopped_sessions.discard(session_id)

            upsert_session(session_id, derive_title(display_text), work_dir)
            if remote_became_ready:
                set_session_remote_state(session_id, remote_session_id, True)
        yield f"data: {json.dumps({'type': 'done'})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.post("/api/chat/stop/{session_id}")
async def stop_chat(session_id: str):
    process = _running_processes.get(session_id)
    if process is None:
        raise HTTPException(status_code=404, detail="no running process for this session")
    _stopped_sessions.add(session_id)
    await _terminate_process(process)
    stop_event = {"type": "error", "message": "用户中止", "ts": time.time()}
    append_event(session_id, stop_event)
    return {"ok": True}


@app.post("/api/sessions/{session_id}/prepare-fork")
async def prepare_fork(session_id: str, req: ForkRequest):
    if session_id in _compacting_sessions:
        raise HTTPException(status_code=409, detail="session is compacting")
    events = load_events(session_id)
    user_event_positions = [i for i, e in enumerate(events) if e.get("type") == "user_input"]
    if not user_event_positions or req.event_index < 0:
        raise HTTPException(status_code=400, detail="invalid event_index")
    event_index = min(req.event_index, len(user_event_positions) - 1)

    target_pos = user_event_positions[event_index]
    events_before = events[:target_pos]
    original_text = events[target_pos].get("text", "")
    new_text = req.new_text if req.new_text is not None and req.new_text.strip() else original_text

    with db_connect() as conn:
        row = conn.execute("SELECT cwd FROM sessions WHERE id = ?", (session_id,)).fetchone()
    cwd = row["cwd"] if row else os.path.expanduser("~")

    new_id = str(uuid.uuid4())
    upsert_session(new_id, derive_title(new_text), cwd)

    with db_connect() as conn:
        conn.execute(
            "UPDATE sessions SET tags = ? WHERE id = ?",
            (f"forked-from-{session_id[:8]}", new_id),
        )

    context = format_context_snippet(events_before)
    if context:
        packed_message = (
            "【以下是之前的对话历史，仅作为参考上下文（不要重复回应历史问题）】\n"
            f"{context}\n\n"
            "【请基于以上历史上下文，回应这个新问题】\n"
            f"{new_text}"
        )
    else:
        packed_message = new_text

    return {
        "session_id": new_id,
        "cwd": cwd,
        "sent_message": packed_message,
        "display_message": new_text,
        "forked_from": session_id,
    }


@app.post("/api/sessions/{session_id}/prepare-inline-edit")
async def prepare_inline_edit(session_id: str, req: ForkRequest):
    if session_id in _running_processes or session_id in _compacting_sessions:
        raise HTTPException(status_code=409, detail="session is busy")

    events = load_events(session_id)
    user_event_positions = [i for i, e in enumerate(events) if e.get("type") == "user_input"]
    if req.event_index < 0 or req.event_index >= len(user_event_positions):
        raise HTTPException(status_code=400, detail="invalid event_index")

    target_pos = user_event_positions[req.event_index]
    events_before = events[:target_pos]
    original_event = events[target_pos]
    original_text = original_event.get("text", "")
    original_images = original_event.get("images", []) or []
    new_text = req.new_text if req.new_text is not None and req.new_text.strip() else original_text

    with db_connect() as conn:
        row = conn.execute("SELECT cwd FROM sessions WHERE id = ?", (session_id,)).fetchone()
    cwd = row["cwd"] if row else os.path.expanduser("~")

    save_events(session_id, events_before)
    upsert_session(session_id, derive_title(new_text), cwd)
    set_session_remote_state(session_id, str(uuid.uuid4()), False)

    context = format_context_snippet(events_before)
    if context:
        packed_message = (
            "【以下是之前的对话历史，仅作为参考上下文（不要重复回应历史问题）】\n"
            f"{context}\n\n"
            "【请基于以上历史上下文，继续这个对话，并回应下面这条经过编辑的新消息】\n"
            f"{new_text}"
        )
    else:
        packed_message = new_text

    return {
        "session_id": session_id,
        "cwd": cwd,
        "sent_message": packed_message,
        "display_message": new_text,
        "images": original_images,
    }


@app.post("/api/sessions/{session_id}/restore-checkpoint")
async def restore_checkpoint(session_id: str, req: RestoreRequest):
    events = load_events(session_id)
    user_event_positions = [i for i, e in enumerate(events) if e.get("type") == "user_input"]
    if req.event_index < 0 or req.event_index >= len(user_event_positions):
        raise HTTPException(status_code=400, detail="invalid event_index")
    ev = events[user_event_positions[req.event_index]]
    cp = ev.get("checkpoint")
    if not cp:
        raise HTTPException(status_code=400, detail="no checkpoint on this turn")

    with db_connect() as conn:
        row = conn.execute("SELECT cwd FROM sessions WHERE id = ?", (session_id,)).fetchone()
    cwd = row["cwd"] if row else ""
    if not cwd:
        raise HTTPException(status_code=400, detail="session has no cwd")

    ok = await restore_git_checkpoint(cwd, cp)
    if not ok:
        raise HTTPException(status_code=500, detail="restore failed")
    return {"ok": True, "cwd": cwd}


@app.post("/api/upload")
async def upload(file: UploadFile = File(...)):
    if file.filename is None:
        raise HTTPException(status_code=400, detail="filename missing")
    ext = Path(file.filename).suffix.lower()
    if ext not in IMAGE_EXTS:
        raise HTTPException(status_code=400, detail=f"unsupported type {ext}")

    data = await file.read()
    if len(data) > MAX_UPLOAD_MB * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"file exceeds {MAX_UPLOAD_MB} MB")

    name = f"{uuid.uuid4().hex}{ext}"
    path = UPLOADS_DIR / name
    path.write_bytes(data)

    return {
        "path": str(path.absolute()),
        "url": f"/uploads/{name}",
        "name": file.filename,
        "size": len(data),
    }


DOC_MIME_EXTS = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.ms-excel.sheet.macroenabled.12": ".xlsm",
    "application/vnd.ms-excel": ".xls",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/xhtml+xml": ".html",
    "application/javascript": ".js",
    "application/json": ".json",
    "application/xml": ".xml",
    "image/svg+xml": ".svg",
    "text/csv": ".csv",
    "text/css": ".css",
    "text/html": ".html",
    "text/javascript": ".js",
    "text/tab-separated-values": ".tsv",
    "text/markdown": ".md",
    "text/plain": ".txt",
    "text/xml": ".xml",
}
MAX_DOC_MB = 30
# Soft cap kept for UI display; we no longer hard-truncate the document text on
# upload. Anything beyond this just gets a "large document" hint in the response.
LARGE_DOC_CHARS_HINT = 200_000
# Argv length safety margin. macOS allows ~256KB total argv; once the prompt
# (UTF-8 bytes) crosses this we route through stdin to avoid E2BIG.
ARGV_STDIN_THRESHOLD = 60_000


def _extract_pdf_text(path: Path) -> str:
    """Extract PDF text. Prefers pdfplumber (better tables/layout) when available,
    falls back to pypdf on any failure (import miss, malformed PDF, table extraction error).
    Each page is prefixed with [Page N] so the model can cite."""
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        pdfplumber = None  # type: ignore
    if pdfplumber is not None:
        try:
            parts: List[str] = []
            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text() or ""
                    tables = []
                    try:
                        for table in page.extract_tables() or []:
                            if not table:
                                continue
                            rows = [" | ".join((cell or "").strip() for cell in row) for row in table]
                            tables.append("\n".join(rows))
                    except Exception:
                        pass
                    section = f"[Page {i}]\n{page_text}"
                    if tables:
                        section += "\n\n" + "\n\n".join(tables)
                    parts.append(section)
            return "\n\n".join(parts)
        except Exception:
            # Any pdfplumber failure (malformed PDF, missing deps, parse error) → fall through to pypdf.
            pass
    import pypdf
    reader = pypdf.PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        parts.append(f"[Page {i}]\n{text}")
    return "\n\n".join(parts)


def _docx_table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        cells = [(cell.text or "").strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if not rows:
        return ""
    if len(rows) == 1:
        return rows[0]
    header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
    return rows[0] + "\n" + header_sep + "\n" + "\n".join(rows[1:])


def _docx_hf_lines(hdr_or_ftr, label: str) -> List[str]:
    """Collect paragraphs and tables from a header/footer container as labeled lines."""
    if hdr_or_ftr is None:
        return []
    out: List[str] = []
    for p in hdr_or_ftr.paragraphs:
        if p.text and p.text.strip():
            out.append(f"[{label}] {p.text}")
    for t in getattr(hdr_or_ftr, "tables", []) or []:
        md = _docx_table_to_markdown(t)
        if md:
            out.append(f"[{label} table]\n{md}")
    return out


def _extract_docx_text(path: Path) -> str:
    """Extract DOCX content preserving the original paragraph/table order.
    Walks the body XML in document order so a 'paragraph → table → paragraph' layout
    survives instead of becoming 'all paragraphs then all tables'.
    Includes default / first-page / even-page headers and footers, plus any
    tables embedded inside them."""
    import docx
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph
    from docx.table import Table
    doc = docx.Document(str(path))
    parts: List[str] = []
    for section in doc.sections:
        parts += _docx_hf_lines(section.header, "Header")
        parts += _docx_hf_lines(section.first_page_header, "Header (first page)")
        parts += _docx_hf_lines(getattr(section, "even_page_header", None), "Header (even page)")
    body = doc.element.body
    para_tag = qn("w:p")
    table_tag = qn("w:tbl")
    for child in body.iterchildren():
        if child.tag == para_tag:
            p = Paragraph(child, doc)
            if p.text:
                parts.append(p.text)
        elif child.tag == table_tag:
            t = Table(child, doc)
            md = _docx_table_to_markdown(t)
            if md:
                parts.append(md)
    for section in doc.sections:
        parts += _docx_hf_lines(section.footer, "Footer")
        parts += _docx_hf_lines(section.first_page_footer, "Footer (first page)")
        parts += _docx_hf_lines(getattr(section, "even_page_footer", None), "Footer (even page)")
    return "\n".join(parts)


def _extract_xlsx_text(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts: List[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(parts)


def _extract_xls_text(path: Path) -> str:
    import xlrd
    wb = xlrd.open_workbook(str(path), on_demand=True)
    try:
        parts: List[str] = []
        for sheet in wb.sheets():
            parts.append(f"[Sheet: {sheet.name}]")
            for row_idx in range(sheet.nrows):
                values = []
                for cell in sheet.row(row_idx):
                    value = cell.value
                    if isinstance(value, float) and value.is_integer():
                        value = int(value)
                    values.append(str(value) if value != "" else "")
                if any(values):
                    parts.append(" | ".join(values))
    finally:
        wb.release_resources()
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    """Extract PowerPoint slides as plain text. Each slide gets a [Slide N] header
    so the model can cite. Pulls title, body text from every shape (recursing into
    grouped shapes), embedded tables (markdown-ified), and the speaker notes pane."""
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

    def walk(shape, title_shape, body_lines: List[str]) -> None:
        # Recurse into groups so text/tables nested inside a group aren't lost.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                walk(child, title_shape, body_lines)
            return
        if title_shape is not None and shape == title_shape:
            return
        if shape.has_text_frame:
            text = (shape.text_frame.text or "").strip()
            if text:
                body_lines.append(text)
        elif shape.has_table:
            rows = []
            for row in shape.table.rows:
                cells = [(c.text or "").strip().replace("\n", " ") for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                if len(rows) > 1:
                    sep = "| " + " | ".join("---" for _ in shape.table.rows[0].cells) + " |"
                    body_lines.append(rows[0] + "\n" + sep + "\n" + "\n".join(rows[1:]))
                else:
                    body_lines.append(rows[0])

    prs = Presentation(str(path))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title_shape = None
        title = ""
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and title_shape.has_text_frame:
                title = (title_shape.text_frame.text or "").strip()
        except Exception:
            title_shape = None
        header = f"[Slide {i}]" + (f" {title}" if title else "")
        body_lines: List[str] = []
        for shape in slide.shapes:
            try:
                walk(shape, title_shape, body_lines)
            except Exception:
                continue  # don't let one bad shape kill the whole slide
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            pass
        section = header
        if body_lines:
            section += "\n" + "\n".join(body_lines)
        if notes_text:
            section += f"\n[Notes] {notes_text}"
        parts.append(section)
    return "\n\n".join(parts)


def _looks_binary(data: bytes) -> bool:
    sample = data[:8192]
    if not sample:
        return False
    if b"\x00" in sample:
        return True
    allowed_controls = {9, 10, 12, 13}
    control_count = sum(1 for b in sample if b < 32 and b not in allowed_controls)
    return control_count / len(sample) > 0.30


def _reject_mojibake(text: str) -> str:
    if not text:
        return text
    replacement_count = text.count("\ufffd")
    if replacement_count and replacement_count / len(text) > 0.01:
        raise HTTPException(status_code=400, detail="unsupported binary file")
    return text


def _decode_text_upload(data: bytes) -> str:
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        return _reject_mojibake(data.decode("utf-16", errors="replace"))
    if data.startswith(b"\xef\xbb\xbf"):
        return _reject_mojibake(data.decode("utf-8-sig", errors="replace"))
    if _looks_binary(data):
        raise HTTPException(status_code=400, detail="unsupported binary file")

    for encoding in ("utf-8", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return _reject_mojibake(data.decode("utf-8", errors="replace"))


def _doc_ext_from_upload(file: UploadFile) -> tuple[str, str]:
    filename = file.filename or "clipboard-file"
    ext = Path(filename).suffix.lower()
    if not ext:
        content_type = (file.content_type or "").split(";", 1)[0].strip().lower()
        ext = DOC_MIME_EXTS.get(content_type, "")
    return ext, filename


@app.post("/api/upload-doc")
async def upload_doc(file: UploadFile = File(...)):
    ext, filename = _doc_ext_from_upload(file)
    data = await file.read()
    if len(data) > MAX_DOC_MB * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"file exceeds {MAX_DOC_MB} MB")

    name = f"{uuid.uuid4().hex}{ext}"
    path = UPLOADS_DIR / name
    path.write_bytes(data)

    try:
        if ext == ".pdf":
            text = _extract_pdf_text(path)
        elif ext == ".docx":
            text = _extract_docx_text(path)
        elif ext == ".pptx":
            text = _extract_pptx_text(path)
        elif ext in (".xlsx", ".xlsm"):
            text = _extract_xlsx_text(path)
        elif ext == ".xls":
            text = _extract_xls_text(path)
        elif ext in (".html", ".htm", ".xhtml"):
            text = _extract_html_text(_decode_text_upload(data))
        else:
            text = _decode_text_upload(data)
    except HTTPException:
        path.unlink(missing_ok=True)
        raise
    except Exception as e:
        path.unlink(missing_ok=True)
        raise HTTPException(status_code=500, detail=f"extract failed: {e}")

    text = text.strip()
    is_large = len(text) > LARGE_DOC_CHARS_HINT

    return {
        "path": str(path.absolute()),
        "name": filename,
        "size": len(data),
        "ext": ext,
        "content": text,
        "length": len(text),
        "truncated": False,
        "large": is_large,
    }


@app.get("/api/doc-content")
async def doc_content(path: str = Query(...)):
    """Read back the extracted text for a doc badge preview.
    Locked to files inside UPLOADS_DIR to prevent path traversal."""
    try:
        target = Path(path).resolve()
    except (OSError, ValueError):
        raise HTTPException(status_code=400, detail="invalid path")
    uploads_root = UPLOADS_DIR.resolve()
    try:
        target.relative_to(uploads_root)
    except ValueError:
        raise HTTPException(status_code=403, detail="path outside uploads directory")
    if not target.is_file():
        raise HTTPException(status_code=404, detail="file not found")

    ext = target.suffix.lower()
    try:
        if ext == ".pdf":
            text = _extract_pdf_text(target)
        elif ext == ".docx":
            text = _extract_docx_text(target)
        elif ext == ".pptx":
            text = _extract_pptx_text(target)
        elif ext in (".xlsx", ".xlsm"):
            text = _extract_xlsx_text(target)
        elif ext == ".xls":
            text = _extract_xls_text(target)
        elif ext in (".html", ".htm", ".xhtml"):
            text = _extract_html_text(_decode_text_upload(target.read_bytes()))
        else:
            text = _decode_text_upload(target.read_bytes())
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"read failed: {e}")
    return {"content": text.strip(), "length": len(text)}


class ExecCodeRequest(BaseModel):
    language: str
    code: str
    timeout: Optional[int] = 10


EXEC_LANG_MAP: Dict[str, List[str]] = {
    "python": ["python3", "-c"],
    "python3": ["python3", "-c"],
    "py": ["python3", "-c"],
    "javascript": ["node", "-e"],
    "js": ["node", "-e"],
    "node": ["node", "-e"],
    "bash": ["bash", "-c"],
    "sh": ["bash", "-c"],
    "shell": ["bash", "-c"],
}


@app.post("/api/exec-code")
async def exec_code(req: ExecCodeRequest):
    lang = (req.language or "").lower().strip()
    cmd = EXEC_LANG_MAP.get(lang)
    if cmd is None:
        raise HTTPException(status_code=400, detail=f"unsupported language: {lang}")
    if not req.code or len(req.code) > 100_000:
        raise HTTPException(status_code=400, detail="code empty or too large")

    timeout = max(1, min(int(req.timeout or 10), 30))
    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd, req.code,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=str(UPLOADS_DIR),
        )
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
        except asyncio.TimeoutError:
            proc.kill()
            return {"stdout": "", "stderr": f"execution timed out after {timeout}s", "returncode": -1, "timed_out": True}
    except FileNotFoundError as e:
        raise HTTPException(status_code=500, detail=f"interpreter not found: {e}")

    return {
        "stdout": stdout.decode("utf-8", errors="replace")[:50_000],
        "stderr": stderr.decode("utf-8", errors="replace")[:10_000],
        "returncode": proc.returncode,
        "timed_out": False,
    }


def _row_to_session(r: sqlite3.Row) -> dict:
    tags = [t for t in (r["tags"] or "").split(",") if t]
    return {
        "id": r["id"],
        "title": r["title"] or "未命名会话",
        "cwd": r["cwd"],
        "created_at": r["created_at"],
        "updated_at": r["updated_at"],
        "pinned": bool(r["pinned"]),
        "archived": bool(r["archived"]),
        "tags": tags,
    }


@app.get("/api/sessions")
async def list_sessions(q: Optional[str] = None, archived: bool = False, tag: Optional[str] = None):
    with db_connect() as conn:
        where = "archived = 1" if archived else "archived = 0"
        rows = conn.execute(
            f"SELECT id, title, cwd, created_at, updated_at, pinned, archived, tags, summary_cache FROM sessions "
            f"WHERE {where} ORDER BY pinned DESC, updated_at DESC LIMIT 500"
        ).fetchall()

    items = []
    for r in rows:
        item = _row_to_session(r)
        item["_summary_cache"] = r["summary_cache"]
        items.append(item)

    if tag:
        items = [i for i in items if tag in i["tags"]]

    if q:
        q_lower = q.lower()
        filtered: List[dict] = []
        for item in items:
            if q_lower in item["title"].lower() or q_lower in ",".join(item["tags"]).lower():
                filtered.append(item)
                continue
            try:
                content = ensure_session_summary_cache(item["id"], item.get("_summary_cache")).lower()
                if q_lower in content:
                    filtered.append(item)
                    continue
                if len(content) >= _SUMMARY_CACHE_LIMIT:
                    full_content = summarize_text_from_events(load_events(item["id"])).lower()
                    if q_lower in full_content:
                        filtered.append(item)
            except Exception:
                continue
        items = filtered

    for item in items:
        item.pop("_summary_cache", None)
    return items


@app.get("/api/sessions/search")
async def search_sessions(q: str = Query(default=""), limit: int = Query(default=10, ge=1, le=30)):
    q_like = f"%{q.strip()}%"
    with db_connect() as conn:
        if q.strip():
            rows = conn.execute(
                """
                SELECT id, title, cwd, updated_at
                FROM sessions
                WHERE archived = 0 AND (title LIKE ? OR cwd LIKE ? OR id LIKE ?)
                ORDER BY pinned DESC, updated_at DESC
                LIMIT ?
                """,
                (q_like, q_like, q_like, limit),
            ).fetchall()
        else:
            rows = conn.execute(
                """
                SELECT id, title, cwd, updated_at
                FROM sessions
                WHERE archived = 0
                ORDER BY pinned DESC, updated_at DESC
                LIMIT ?
                """,
                (limit,),
            ).fetchall()
    return [dict(r) for r in rows]


@app.get("/api/sessions/{session_id}")
async def get_session(session_id: str):
    with db_connect() as conn:
        row = conn.execute(
            "SELECT id, title, cwd, created_at, updated_at, pinned, archived, tags FROM sessions WHERE id = ?",
            (session_id,),
        ).fetchone()
    if row is None:
        raise HTTPException(status_code=404, detail="session not found")
    data = _row_to_session(row)
    data["events"] = load_events(session_id)
    data["compact_backups"] = [
        {"name": p.name, "created_at": p.stat().st_mtime, "size": p.stat().st_size}
        for p in sorted(iter_session_compact_backups(session_id), key=lambda x: x.stat().st_mtime, reverse=True)
    ]
    return data


@app.patch("/api/sessions/{session_id}")
async def patch_session(session_id: str, req: SessionPatch):
    updates: List[str] = []
    params: List = []
    if req.title is not None:
        updates += ["title = ?", "manual_title = 1"]
        params.append(req.title)
    if req.pinned is not None:
        updates.append("pinned = ?")
        params.append(1 if req.pinned else 0)
    if req.archived is not None:
        updates.append("archived = ?")
        params.append(1 if req.archived else 0)
    if req.tags is not None:
        updates.append("tags = ?")
        params.append(req.tags)
    if not updates:
        return {"ok": True}
    params.append(session_id)
    with db_connect() as conn:
        conn.execute(f"UPDATE sessions SET {', '.join(updates)} WHERE id = ?", params)
    return {"ok": True}


@app.delete("/api/sessions/{session_id}")
async def delete_session(session_id: str):
    if session_id in _running_processes or session_id in _compacting_sessions:
        raise HTTPException(status_code=409, detail="session is busy")
    with db_connect() as conn:
        conn.execute("DELETE FROM sessions WHERE id = ?", (session_id,))
        conn.execute("DELETE FROM session_usage WHERE session_id = ?", (session_id,))
        conn.execute("DELETE FROM tool_calls WHERE session_id = ?", (session_id,))
        conn.execute("DELETE FROM memories WHERE scope = ?", (f"session:{session_id}",))
    path = HISTORY_DIR / f"{session_id}.jsonl"
    if path.exists():
        path.unlink()
    for backup in iter_session_compact_backups(session_id):
        backup.unlink(missing_ok=True)
    return {"ok": True}


@app.post("/api/sessions/{session_id}/clear")
async def clear_session(session_id: str):
    if session_id in _running_processes or session_id in _compacting_sessions:
        raise HTTPException(status_code=409, detail="session is busy")
    save_events(session_id, [])
    with db_connect() as conn:
        conn.execute("UPDATE sessions SET title = '新会话', manual_title = 0, updated_at = ? WHERE id = ?", (time.time(), session_id))
    set_session_remote_state(session_id, "", False)
    return {"ok": True}


@app.post("/api/sessions/{session_id}/compact")
async def compact_session(session_id: str, keep_last: int = Query(default=2, ge=1, le=10)):
    if session_id in _running_processes or session_id in _compacting_sessions:
        raise HTTPException(status_code=409, detail="session is busy")
    _compacting_sessions.add(session_id)
    try:
        events = load_events(session_id)
        if len(events) < 4:
            return {"ok": True, "skipped": True, "reason": "history too short"}

        user_indices = [i for i, e in enumerate(events) if e.get("type") == "user_input"]
        if len(user_indices) <= keep_last:
            return {"ok": True, "skipped": True, "reason": "history too short"}

        split_at = user_indices[-keep_last]
        old_events, new_events = events[:split_at], events[split_at:]
        snippet = format_context_snippet(old_events, max_chars=12000)
        summary_prompt = (
            "请把以下对话历史压缩成一份延续工作所需的精简摘要，"
            "覆盖：目标、关键决策、已修改文件、未完成工作、风险与约定。"
            "用 markdown 列表，不超过 30 行。\n\n"
            + snippet
        )
        proc = await asyncio.create_subprocess_exec(
            *claude_cli_argv("-p", summary_prompt, "--output-format", "text", "--model", "haiku"),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=60)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            raise HTTPException(status_code=504, detail="summary timeout")

        summary = stdout.decode("utf-8", errors="replace").strip()
        if not summary:
            raise HTTPException(status_code=500, detail="empty summary")

        src = HISTORY_DIR / f"{session_id}.jsonl"
        backup_name = ""
        if src.exists():
            backup = HISTORY_DIR / f"{session_id}.before-compact-{int(time.time())}.jsonl"
            backup.write_bytes(src.read_bytes())
            backup_name = backup.name
            prune_session_compact_backups(session_id)

        compacted = [
            {
                "type": "user_input",
                "text": f"【会话已压缩 · 以下为之前对话的摘要】\n\n{summary}",
                "ts": time.time(),
                "compacted": True,
            }
        ] + new_events
        save_events(session_id, compacted)
        set_session_remote_state(session_id, "", False)
        return {"ok": True, "kept_turns": keep_last, "backup": backup_name}
    except ClaudeCliResolutionError as e:
        raise HTTPException(status_code=500, detail=str(e))
    except FileNotFoundError:
        raise HTTPException(status_code=500, detail="claude CLI not found in PATH")
    finally:
        _compacting_sessions.discard(session_id)


@app.post("/api/sessions/{session_id}/suggest-title")
async def suggest_title(session_id: str):
    events = load_events(session_id)
    if not events:
        raise HTTPException(status_code=404, detail="empty session")
    summary = summarize_text_from_events(events)[:3000]
    if not summary.strip():
        raise HTTPException(status_code=400, detail="no textual content")
    prompt = f"根据下面的对话，用中文生成一个不超过15字、不带引号的会话标题（只输出标题本身）：\n\n{summary}"
    try:
        proc = await asyncio.create_subprocess_exec(
            *claude_cli_argv("-p", prompt, "--output-format", "text"),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=60)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            raise HTTPException(status_code=504, detail="title generation timeout")
    except HTTPException:
        raise
    except ClaudeCliResolutionError as e:
        raise HTTPException(status_code=500, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"title generation failed: {e}")
    title = stdout.decode("utf-8", errors="replace").strip().splitlines()[0].strip(' "\'"""''').strip()[:60]
    if not title:
        raise HTTPException(status_code=500, detail="empty title")
    with db_connect() as conn:
        conn.execute("UPDATE sessions SET title = ?, manual_title = 1 WHERE id = ?", (title, session_id))
    return {"title": title}


@app.get("/api/sessions/{session_id}/export")
async def export_session(session_id: str):
    events = load_events(session_id)
    if not events:
        raise HTTPException(status_code=404, detail="session not found")
    with db_connect() as conn:
        row = conn.execute("SELECT title FROM sessions WHERE id = ?", (session_id,)).fetchone()
    title = row["title"] if row and row["title"] else session_id

    lines: List[str] = [f"# {title}", "", f"_会话 ID: {session_id}_", ""]
    for ev in events:
        t = ev.get("type")
        if t == "user_input":
            lines += ["## 👤 用户", "", ev.get("text", "")]
            for img in ev.get("images", []) or []:
                lines.append(f"![image]({img})")
            lines.append("")
        elif t == "assistant":
            content = (ev.get("message") or {}).get("content") or []
            for block in content:
                if block.get("type") == "text":
                    lines += ["## 🤖 Claude", "", block.get("text", ""), ""]
                elif block.get("type") == "tool_use":
                    name = block.get("name", "?")
                    lines += [f"### 🔧 工具调用: `{name}`", "", "```json",
                              json.dumps(block.get("input", {}), ensure_ascii=False, indent=2), "```", ""]
        elif t == "user":
            content = (ev.get("message") or {}).get("content") or []
            for block in content:
                if block.get("type") == "tool_result":
                    ct = block.get("content", "")
                    if isinstance(ct, list):
                        ct = "\n".join(c.get("text", "") if isinstance(c, dict) else str(c) for c in ct)
                    lines += ["### 📋 工具结果", "", "```", str(ct)[:5000], "```", ""]

    md = "\n".join(lines)
    return Response(
        content=md,
        media_type="text/markdown; charset=utf-8",
        headers={"Content-Disposition": f'attachment; filename="{session_id}.md"'},
    )


@app.get("/api/sessions/{session_id}/usage")
async def get_session_usage(session_id: str, limit: int = Query(default=20, ge=1, le=100)):
    with db_connect() as conn:
        total = conn.execute(
            """
            SELECT
                COALESCE(SUM(input_tokens), 0) AS input_tokens,
                COALESCE(SUM(output_tokens), 0) AS output_tokens,
                COALESCE(SUM(cache_read_input_tokens), 0) AS cache_read_input_tokens,
                COALESCE(SUM(cache_creation_input_tokens), 0) AS cache_creation_input_tokens,
                COALESCE(SUM(total_cost_usd), 0) AS total_cost_usd
            FROM session_usage
            WHERE session_id = ?
            """,
            (session_id,),
        ).fetchone()
        rows = conn.execute(
            """
            SELECT turn_idx, input_tokens, output_tokens, cache_read_input_tokens,
                   cache_creation_input_tokens, total_cost_usd, ts
            FROM session_usage
            WHERE session_id = ?
            ORDER BY turn_idx DESC
            LIMIT ?
            """,
            (session_id, limit),
        ).fetchall()
    return {
        "session_id": session_id,
        "total": dict(total) if total else {},
        "recent": [dict(r) for r in rows],
    }


@app.get("/api/sessions/{session_id}/mention")
async def mention_session(session_id: str, max_chars: int = Query(default=5000, ge=500, le=12000)):
    events = load_events(session_id)
    if not events:
        raise HTTPException(status_code=404, detail="session not found")
    text = summarize_text_from_events(events)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if len(text) > max_chars:
        head = text[: max_chars // 2].rstrip()
        tail = text[-(max_chars // 2) :].lstrip()
        text = head + "\n\n...[session content truncated]...\n\n" + tail
    with db_connect() as conn:
        row = conn.execute("SELECT title FROM sessions WHERE id = ?", (session_id,)).fetchone()
    title = row["title"] if row and row["title"] else session_id[:8]
    return {"id": session_id, "title": title, "content": f"Referenced session: {title}\n\n{text}"}


@app.post("/api/projects/scan")
async def scan_project(cwd: str = Query(...)):
    project_dir = Path(os.path.expanduser(cwd)).resolve()
    if not project_dir.is_dir():
        raise HTTPException(status_code=400, detail="cwd not found")

    probes = [
        "README.md", "README", "package.json", "pyproject.toml", "Cargo.toml",
        "go.mod", "tsconfig.json", "Makefile", ".gitignore",
    ]
    snippets: List[str] = []
    for name in probes:
        p = project_dir / name
        try:
            if p.is_file() and p.stat().st_size < 32_000:
                snippets.append(f"--- {name} ---\n" + p.read_text(encoding="utf-8", errors="replace"))
        except OSError:
            continue

    scan_ignored_dirs = IGNORED_DIRS | {"history", "uploads", "dist", "build", ".pycache_check"}
    try:
        tree_lines: List[str] = []
        for entry in sorted(project_dir.iterdir(), key=lambda x: x.name)[:50]:
            if entry.name.startswith("."):
                continue
            if entry.is_dir():
                if entry.name in scan_ignored_dirs:
                    continue
                tree_lines.append(f"DIR {entry.name}/")
                try:
                    for sub in sorted(entry.iterdir(), key=lambda x: x.name)[:20]:
                        if not sub.name.startswith("."):
                            tree_lines.append(f"   {sub.name}{'/' if sub.is_dir() else ''}")
                except OSError:
                    pass
            else:
                tree_lines.append(f"FILE {entry.name}")
        snippets.append("--- directory ---\n" + "\n".join(tree_lines))
    except OSError:
        pass

    return {"cwd": str(project_dir), "context": "\n\n".join(snippets)[:20_000]}


@app.get("/api/memories")
async def list_memories(scope: Optional[str] = None, q: str = Query(default="")):
    clauses = []
    params: List[object] = []
    if scope:
        clauses.append("scope = ?")
        params.append(scope)
    if q.strip():
        clauses.append("content LIKE ?")
        params.append(f"%{q.strip()}%")
    where = ("WHERE " + " AND ".join(clauses)) if clauses else ""
    with db_connect() as conn:
        rows = conn.execute(
            f"""
            SELECT id, content, enabled, scope, created_at, updated_at
            FROM memories
            {where}
            ORDER BY updated_at DESC
            LIMIT 100
            """,
            params,
        ).fetchall()
    return [dict(r) for r in rows]


@app.get("/api/memories/active")
async def active_memories(cwd: str = Query(default=""), session_id: str = Query(default="")):
    return load_enabled_memories(cwd, session_id)


@app.post("/api/memories")
async def create_memory(req: MemoryRequest):
    mid = uuid.uuid4().hex
    now = time.time()
    content = req.content.strip()
    if not content:
        raise HTTPException(status_code=400, detail="content required")
    scope = normalize_memory_scope(req.scope)
    with db_connect() as conn:
        conn.execute(
            "INSERT INTO memories (id, content, enabled, scope, created_at, updated_at) VALUES (?, ?, ?, ?, ?, ?)",
            (mid, content, 1 if req.enabled else 0, scope, now, now),
        )
    return {"id": mid}


@app.put("/api/memories/{memory_id}")
async def update_memory(memory_id: str, req: MemoryRequest):
    content = req.content.strip()
    if not content:
        raise HTTPException(status_code=400, detail="content required")
    scope = normalize_memory_scope(req.scope)
    with db_connect() as conn:
        cursor = conn.execute(
            "UPDATE memories SET content = ?, enabled = ?, scope = ?, updated_at = ? WHERE id = ?",
            (content, 1 if req.enabled else 0, scope, time.time(), memory_id),
        )
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="memory not found")
    return {"ok": True}


@app.delete("/api/memories/{memory_id}")
async def delete_memory(memory_id: str):
    with db_connect() as conn:
        cursor = conn.execute("DELETE FROM memories WHERE id = ?", (memory_id,))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="memory not found")
    return {"ok": True}


@app.get("/api/memories/search")
async def search_memories(q: str = Query(default=""), limit: int = Query(default=10, ge=1, le=30)):
    with db_connect() as conn:
        rows = conn.execute(
            """
            SELECT id, content, enabled, scope, updated_at
            FROM memories
            WHERE content LIKE ?
            ORDER BY enabled DESC, updated_at DESC
            LIMIT ?
            """,
            (f"%{q.strip()}%", limit),
        ).fetchall()
    return [dict(r) for r in rows]


@app.get("/api/prompts")
async def list_prompts():
    with db_connect() as conn:
        rows = conn.execute(
            "SELECT id, name, content, slash_trigger, created_at FROM prompts ORDER BY created_at DESC"
        ).fetchall()
    return [
        {"id": r["id"], "name": r["name"], "content": r["content"], "slash_trigger": r["slash_trigger"], "created_at": r["created_at"]}
        for r in rows
    ]


@app.get("/api/prompts/search")
async def search_prompts(q: str = Query(default=""), limit: int = Query(default=10, ge=1, le=30)):
    q_like = f"%{q.strip()}%"
    with db_connect() as conn:
        rows = conn.execute(
            """
            SELECT id, name, content, slash_trigger, created_at
            FROM prompts
            WHERE name LIKE ? OR content LIKE ? OR slash_trigger LIKE ?
            ORDER BY created_at DESC
            LIMIT ?
            """,
            (q_like, q_like, q_like, limit),
        ).fetchall()
    return [dict(r) for r in rows]


@app.post("/api/prompts")
async def create_prompt(req: PromptRequest):
    pid = uuid.uuid4().hex
    with db_connect() as conn:
        conn.execute(
            "INSERT INTO prompts (id, name, content, slash_trigger, created_at) VALUES (?, ?, ?, ?, ?)",
            (pid, req.name, req.content, (req.slash_trigger or "").strip().lstrip("/"), time.time()),
        )
    return {"id": pid}


@app.put("/api/prompts/{prompt_id}")
async def update_prompt(prompt_id: str, req: PromptRequest):
    with db_connect() as conn:
        cursor = conn.execute(
            "UPDATE prompts SET name = ?, content = ?, slash_trigger = ? WHERE id = ?",
            (req.name, req.content, (req.slash_trigger or "").strip().lstrip("/"), prompt_id),
        )
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="prompt not found")
    return {"ok": True}


@app.delete("/api/prompts/{prompt_id}")
async def delete_prompt(prompt_id: str):
    with db_connect() as conn:
        cursor = conn.execute("DELETE FROM prompts WHERE id = ?", (prompt_id,))
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail="prompt not found")
    return {"ok": True}


@app.post("/api/suggest-followups")
async def suggest_followups(session_id: str = ""):
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id required")
    events = load_events(session_id)
    if not events:
        return {"suggestions": []}
    snippet = summarize_text_from_events(events[-20:])[-3000:]
    if not snippet.strip():
        return {"suggestions": []}
    prompt = (
        "根据以下对话内容，生成3个用户可能想继续追问的简短问题（每个不超过20字）。"
        "只输出3行，每行一个问题，不要编号、不要引号、不要其他内容。\n\n"
        f"{snippet}"
    )
    try:
        proc = await asyncio.create_subprocess_exec(
            *claude_cli_argv("-p", prompt, "--output-format", "text", "--model", "haiku"),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return {"suggestions": []}
    except Exception:
        return {"suggestions": []}
    lines = [l.strip() for l in stdout.decode("utf-8", errors="replace").splitlines() if l.strip()]
    suggestions = [l.lstrip("0123456789.-、）) ") for l in lines[:3]]
    return {"suggestions": suggestions}


# ===== MCP Management =====

_CLAUDE_CONFIG_PATH = Path.home() / ".claude.json"
_PROJECT_MCP_FILENAME = ".mcp.json"
_DISABLED_MCP_SERVERS_KEY = "claudeWebDisabledMcpServers"
_MCP_SCOPES = {"local", "user", "project"}


def _read_json_object(path: Path) -> dict:
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"invalid JSON in {path}: {e.msg}")
    except OSError as e:
        raise HTTPException(status_code=500, detail=f"cannot read {path}: {e}")
    if not isinstance(data, dict):
        raise HTTPException(status_code=500, detail=f"invalid JSON object in {path}")
    return data


def _write_json_object(path: Path, data: dict) -> None:
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    except OSError as e:
        raise HTTPException(status_code=500, detail=f"cannot write {path}: {e}")


def _normalize_mcp_scope(scope: Optional[str]) -> str:
    normalized = (scope or "local").strip().lower()
    if normalized not in _MCP_SCOPES:
        raise HTTPException(status_code=400, detail="scope must be local, user, or project")
    return normalized


def _resolve_mcp_cwd(cwd: Optional[str]) -> Path:
    raw = (cwd or "").strip() or "~"
    target = Path(os.path.expanduser(raw)).resolve()
    if not target.is_dir():
        raise HTTPException(status_code=400, detail=f"invalid cwd: {raw}")
    return target


def _dict_value(parent: dict, key: str) -> dict:
    value = parent.get(key)
    if not isinstance(value, dict):
        value = {}
        parent[key] = value
    return value


def _source_label(scope: str) -> str:
    return {
        "local": "Local",
        "user": "User",
        "project": "Project",
    }.get(scope, scope)


def _mcp_sources(cwd: Optional[str], create: bool = False) -> list[dict]:
    project_dir = _resolve_mcp_cwd(cwd)
    project_key = str(project_dir)
    claude_data = _read_json_object(_CLAUDE_CONFIG_PATH)
    project_mcp_path = project_dir / _PROJECT_MCP_FILENAME
    project_mcp_data = _read_json_object(project_mcp_path)

    projects = _dict_value(claude_data, "projects") if create else claude_data.get("projects", {})
    if not isinstance(projects, dict):
        projects = {}
    local_project = projects.get(project_key)
    if create:
        local_project = projects.setdefault(project_key, {})
    if not isinstance(local_project, dict):
        local_project = {}

    project_choice = local_project if isinstance(local_project, dict) else {}
    disabled_project_servers = project_choice.get("disabledMcpjsonServers", [])
    if not isinstance(disabled_project_servers, list):
        disabled_project_servers = []

    return [
        {
            "scope": "user",
            "path": _CLAUDE_CONFIG_PATH,
            "data": claude_data,
            "servers": claude_data.get("mcpServers", {}) if isinstance(claude_data.get("mcpServers", {}), dict) else {},
            "disabled_servers": claude_data.get(_DISABLED_MCP_SERVERS_KEY, {}) if isinstance(claude_data.get(_DISABLED_MCP_SERVERS_KEY, {}), dict) else {},
        },
        {
            "scope": "local",
            "path": _CLAUDE_CONFIG_PATH,
            "data": claude_data,
            "servers": local_project.get("mcpServers", {}) if isinstance(local_project.get("mcpServers", {}), dict) else {},
            "disabled_servers": local_project.get(_DISABLED_MCP_SERVERS_KEY, {}) if isinstance(local_project.get(_DISABLED_MCP_SERVERS_KEY, {}), dict) else {},
            "project_key": project_key,
            "project": local_project,
        },
        {
            "scope": "project",
            "path": project_mcp_path,
            "data": project_mcp_data,
            "servers": project_mcp_data.get("mcpServers", {}) if isinstance(project_mcp_data.get("mcpServers", {}), dict) else {},
            "disabled_names": set(str(v) for v in disabled_project_servers),
            "claude_data": claude_data,
            "project_key": project_key,
            "project": local_project,
        },
    ]


def _mcp_target(scope: str, cwd: Optional[str]) -> dict:
    normalized = _normalize_mcp_scope(scope)
    sources = _mcp_sources(cwd, create=True)
    for source in sources:
        if source["scope"] == normalized:
            if normalized == "user":
                source["servers"] = _dict_value(source["data"], "mcpServers")
                source["disabled_servers"] = _dict_value(source["data"], _DISABLED_MCP_SERVERS_KEY)
            elif normalized == "local":
                source["servers"] = _dict_value(source["project"], "mcpServers")
                source["disabled_servers"] = _dict_value(source["project"], _DISABLED_MCP_SERVERS_KEY)
            else:
                source["servers"] = _dict_value(source["data"], "mcpServers")
                disabled = source["project"].get("disabledMcpjsonServers")
                if not isinstance(disabled, list):
                    disabled = []
                    source["project"]["disabledMcpjsonServers"] = disabled
                source["disabled_names"] = set(str(v) for v in disabled)
            return source
    raise HTTPException(status_code=400, detail="invalid scope")


def _save_mcp_source(source: dict, save_claude_choices: bool = False) -> None:
    _write_json_object(source["path"], source["data"])
    if source["scope"] == "project" and save_claude_choices:
        _write_json_object(_CLAUDE_CONFIG_PATH, source["claude_data"])


def _find_mcp_source(name: str, scope: Optional[str], cwd: Optional[str]) -> dict:
    if scope:
        source = _mcp_target(scope, cwd)
        if _mcp_config_in_source(source, name) is None:
            raise HTTPException(status_code=404, detail=f"server '{name}' not found")
        return source

    matches = []
    for source in _mcp_sources(cwd, create=True):
        if _mcp_config_in_source(source, name) is not None:
            matches.append(source)
    if not matches:
        raise HTTPException(status_code=404, detail=f"server '{name}' not found")
    if len(matches) > 1:
        scopes = ", ".join(_source_label(m["scope"]) for m in matches)
        raise HTTPException(status_code=409, detail=f"server '{name}' exists in multiple scopes: {scopes}")
    return matches[0]


def _mcp_config_in_source(source: dict, name: str) -> Optional[dict]:
    servers = source.get("servers") or {}
    if name in servers and isinstance(servers[name], dict):
        return servers[name]
    disabled = source.get("disabled_servers") or {}
    if name in disabled and isinstance(disabled[name], dict):
        return disabled[name]
    return None


def _is_mcp_disabled(source: dict, name: str) -> bool:
    if source["scope"] == "project":
        return name in (source.get("disabled_names") or set())
    return name in (source.get("disabled_servers") or {})


def _set_mcp_disabled(source: dict, name: str, disabled: bool) -> bool:
    if source["scope"] == "project":
        project = source["project"]
        disabled_list = project.get("disabledMcpjsonServers")
        if not isinstance(disabled_list, list):
            disabled_list = []
            project["disabledMcpjsonServers"] = disabled_list
        if disabled and name not in disabled_list:
            disabled_list.append(name)
        elif not disabled:
            project["disabledMcpjsonServers"] = [n for n in disabled_list if n != name]
        return True

    servers = source["servers"]
    disabled_servers = source["disabled_servers"]
    if disabled:
        cfg = servers.pop(name, disabled_servers.get(name))
        if cfg is not None:
            cfg.pop("disabled", None)
            disabled_servers[name] = cfg
    else:
        cfg = disabled_servers.pop(name, servers.get(name))
        if cfg is not None:
            cfg.pop("disabled", None)
            servers[name] = cfg
    return False


def _mask_mapping(values: Optional[dict]) -> dict:
    if not values:
        return {}
    masked = {}
    for k, v in values.items():
        value = str(v)
        if any(s in k.lower() for s in ("token", "key", "secret", "password", "credential", "auth")):
            masked[k] = value[:4] + "***" if len(value) > 4 else "***"
        else:
            masked[k] = value
    return masked


def _mcp_transport(cfg: dict) -> str:
    transport = str(cfg.get("type") or cfg.get("transport") or "").strip().lower()
    if transport:
        return transport
    if cfg.get("url"):
        return "http"
    return "stdio"


def _format_mcp_server(name: str, cfg: dict, source: dict, disabled: bool) -> dict:
    transport = _mcp_transport(cfg)
    return {
        "name": name,
        "scope": source["scope"],
        "scope_label": _source_label(source["scope"]),
        "config_path": str(source["path"]),
        "type": transport,
        "command": cfg.get("command", ""),
        "args": cfg.get("args", []) if isinstance(cfg.get("args", []), list) else [],
        "url": cfg.get("url", ""),
        "env": _mask_mapping(cfg.get("env") if isinstance(cfg.get("env"), dict) else {}),
        "headers": _mask_mapping(cfg.get("headers") if isinstance(cfg.get("headers"), dict) else {}),
        "disabled": disabled,
    }


def _stdio_config_from_request(req: "McpServerRequest") -> dict:
    command = (req.command or "").strip()
    if not command:
        raise HTTPException(status_code=400, detail="command is required")
    cfg = {
        "type": "stdio",
        "command": command,
        "args": req.args or [],
    }
    if req.env:
        cfg["env"] = req.env
    return cfg


class McpServerRequest(BaseModel):
    command: Optional[str] = None
    args: Optional[List[str]] = None
    env: Optional[Dict[str, str]] = None
    disabled: Optional[bool] = None


class McpServerPatchRequest(BaseModel):
    command: Optional[str] = None
    args: Optional[List[str]] = None
    env: Optional[Dict[str, str]] = None
    disabled: Optional[bool] = None


@app.get("/api/mcp/servers")
async def list_mcp_servers(cwd: Optional[str] = Query(default=None)):
    sources = _mcp_sources(cwd)
    result = []
    for source in sources:
        servers = source.get("servers") or {}
        for name, cfg in servers.items():
            if isinstance(cfg, dict):
                result.append(_format_mcp_server(name, cfg, source, _is_mcp_disabled(source, name)))
        for name, cfg in (source.get("disabled_servers") or {}).items():
            if name not in servers and isinstance(cfg, dict):
                result.append(_format_mcp_server(name, cfg, source, True))
    return {
        "servers": result,
        "cwd": str(_resolve_mcp_cwd(cwd)),
        "config_path": str(_CLAUDE_CONFIG_PATH),
        "config_paths": {
            "user": str(_CLAUDE_CONFIG_PATH),
            "local": str(_CLAUDE_CONFIG_PATH),
            "project": str(_resolve_mcp_cwd(cwd) / _PROJECT_MCP_FILENAME),
        },
    }


@app.post("/api/mcp/servers/{name}")
async def add_mcp_server(
    name: str,
    req: McpServerRequest,
    cwd: Optional[str] = Query(default=None),
    scope: str = Query(default="local"),
):
    target = _mcp_target(scope, cwd)
    if _mcp_config_in_source(target, name) is not None:
        raise HTTPException(status_code=409, detail=f"server '{name}' already exists")
    cfg = _stdio_config_from_request(req)
    if req.disabled:
        if target["scope"] == "project":
            target["servers"][name] = cfg
            save_choices = _set_mcp_disabled(target, name, True)
            _save_mcp_source(target, save_claude_choices=save_choices)
        else:
            target["disabled_servers"][name] = cfg
            _save_mcp_source(target)
    else:
        target["servers"][name] = cfg
        _save_mcp_source(target)
    return {"ok": True}


@app.patch("/api/mcp/servers/{name}")
async def patch_mcp_server(
    name: str,
    req: McpServerPatchRequest,
    cwd: Optional[str] = Query(default=None),
    scope: Optional[str] = Query(default=None),
):
    target = _find_mcp_source(name, scope, cwd)
    cfg = _mcp_config_in_source(target, name)
    if cfg is None:
        raise HTTPException(status_code=404, detail=f"server '{name}' not found")
    if req.command is not None:
        command = req.command.strip()
        if not command:
            raise HTTPException(status_code=400, detail="command is required")
        cfg["command"] = command
        cfg.setdefault("type", "stdio")
    if req.args is not None:
        cfg["args"] = req.args
    if req.env is not None:
        cfg["env"] = req.env
    save_choices = False
    cfg.pop("disabled", None)
    if req.disabled is not None:
        save_choices = _set_mcp_disabled(target, name, req.disabled)
    _save_mcp_source(target, save_claude_choices=save_choices)
    return {"ok": True}


@app.delete("/api/mcp/servers/{name}")
async def delete_mcp_server(
    name: str,
    cwd: Optional[str] = Query(default=None),
    scope: Optional[str] = Query(default=None),
):
    target = _find_mcp_source(name, scope, cwd)
    target["servers"].pop(name, None)
    if target["scope"] == "project":
        save_choices = _set_mcp_disabled(target, name, False)
        _save_mcp_source(target, save_claude_choices=save_choices)
    else:
        target["disabled_servers"].pop(name, None)
        _save_mcp_source(target)
    return {"ok": True}


@app.get("/api/cwds")
async def list_cwds():
    with db_connect() as conn:
        rows = conn.execute(
            "SELECT cwd, MAX(updated_at) AS last FROM sessions WHERE cwd <> '' GROUP BY cwd ORDER BY last DESC LIMIT 10"
        ).fetchall()
    return [r["cwd"] for r in rows]


@app.get("/api/tags")
async def list_tags():
    with db_connect() as conn:
        rows = conn.execute("SELECT tags FROM sessions WHERE tags <> '' AND archived = 0").fetchall()
    counts: Dict[str, int] = defaultdict(int)
    for r in rows:
        for t in (r["tags"] or "").split(","):
            t = t.strip()
            if t:
                counts[t] += 1
    return [{"name": k, "count": v} for k, v in sorted(counts.items(), key=lambda x: -x[1])]


@app.get("/api/stats")
async def stats():
    await ensure_stats_backfilled()
    with db_connect() as conn:
        total_sessions = conn.execute("SELECT COUNT(*) AS c FROM sessions").fetchone()["c"]
        usage = conn.execute(
            """
            SELECT
                COALESCE(SUM(total_cost_usd), 0) AS total_cost,
                COALESCE(SUM(duration_ms), 0) AS total_duration,
                COUNT(*) AS total_turns
            FROM session_usage
            """
        ).fetchone()
        daily_rows = conn.execute(
            """
            SELECT date(ts, 'unixepoch', 'localtime') AS day,
                   COALESCE(SUM(total_cost_usd), 0) AS cost,
                   COUNT(*) AS turns
            FROM session_usage
            GROUP BY day
            ORDER BY day
            """
        ).fetchall()
        tool_rows = conn.execute(
            """
            SELECT tool_name AS name, COUNT(*) AS count
            FROM tool_calls
            GROUP BY tool_name
            ORDER BY count DESC
            LIMIT 10
            """
        ).fetchall()
    return {
        "total_cost_usd": round(float(usage["total_cost"] or 0), 4),
        "total_duration_ms": float(usage["total_duration"] or 0),
        "total_sessions": total_sessions,
        "total_turns": int(usage["total_turns"] or 0),
        "daily": [
            {"date": r["day"], "cost": round(float(r["cost"] or 0), 4), "turns": r["turns"]}
            for r in daily_rows
            if r["day"] is not None
        ],
        "tools": [{"name": r["name"], "count": r["count"]} for r in tool_rows],
    }


async def _list_files_via_git(base: Path, q_lower: str, limit: int) -> Optional[List[dict]]:
    try:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", str(base), "ls-files",
            "--cached", "--others", "--exclude-standard",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.DEVNULL,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=5)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return None
        if proc.returncode != 0:
            return None
    except (FileNotFoundError, asyncio.TimeoutError):
        return None
    except Exception:
        return None

    results: List[dict] = []
    for rel in stdout.decode("utf-8", errors="replace").splitlines():
        rel = rel.strip()
        if not rel:
            continue
        if q_lower and q_lower not in rel.lower():
            continue
        results.append({"path": str(base / rel), "rel": rel})
        if len(results) >= limit:
            break
    return results


@app.get("/api/files")
async def list_files(cwd: str = Query(...), q: str = Query(default=""), limit: int = Query(default=30)):
    base = Path(os.path.expanduser(cwd)).resolve()
    if not base.exists() or not base.is_dir():
        return []
    q_lower = q.lower()

    git_results = await _list_files_via_git(base, q_lower, limit)
    if git_results is not None:
        return git_results

    results: List[dict] = []
    for root, dirs, files in os.walk(str(base)):
        dirs[:] = [d for d in dirs if d not in IGNORED_DIRS and not d.startswith(".")]
        for f in files:
            if f.startswith("."):
                continue
            full = Path(root) / f
            try:
                rel = str(full.relative_to(base))
            except ValueError:
                continue
            if q_lower and q_lower not in rel.lower():
                continue
            results.append({"path": str(full), "rel": rel})
            if len(results) >= limit:
                return results
    return results


@app.get("/api/git")
async def git_status(cwd: str = Query(...)):
    target = os.path.expanduser(cwd)
    if not os.path.isdir(target):
        return {"branch": "", "dirty": 0, "available": False}
    try:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", target, "status", "--porcelain=v1", "--branch",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=5)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return {"branch": "", "dirty": 0, "available": False}
    except Exception:
        return {"branch": "", "dirty": 0, "available": False}
    if proc.returncode != 0:
        return {"branch": "", "dirty": 0, "available": False}
    branch = ""
    dirty = 0
    for line in stdout.decode("utf-8", errors="replace").splitlines():
        if line.startswith("##"):
            header = line[2:].strip()
            branch = header.split("...")[0].strip()
        else:
            dirty += 1
    return {"branch": branch, "dirty": dirty, "available": True}


app.mount("/uploads", StaticFiles(directory=str(UPLOADS_DIR)), name="uploads")


class _TextExtractor(HTMLParser):
    _SKIP_TAGS = {"script", "style", "noscript", "svg", "iframe", "head"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: List[str] = []
        self._skip_depth = 0
        self._table_depth = 0
        self._active_rowspans: Dict[int, int] = {}
        self._new_rowspans: Dict[int, int] = {}
        self._current_row: Optional[List[str]] = None
        self._current_cell: Optional[List[str]] = None
        self._current_colspan = 1
        self._current_rowspan = 1

    def handle_starttag(self, tag: str, attrs: List) -> None:
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1
        elif self._skip_depth > 0:
            return
        elif tag == "table":
            if self._table_depth == 0:
                self._parts.append("\n")
                self._active_rowspans = {}
            self._table_depth += 1
        elif self._table_depth > 0:
            if self._table_depth == 1 and tag == "tr":
                self._start_table_row()
            elif self._table_depth == 1 and tag in {"td", "th"}:
                self._start_table_cell(attrs)
            elif tag == "br" and self._current_cell is not None:
                self._current_cell.append("\n")
        elif tag in {"p", "br", "div", "li", "h1", "h2", "h3", "h4", "h5", "h6", "tr"}:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1
        elif self._skip_depth > 0:
            return
        elif tag in {"td", "th"} and self._table_depth == 1:
            self._end_table_cell()
        elif tag == "tr" and self._table_depth == 1:
            self._end_table_row()
        elif tag == "table" and self._table_depth > 0:
            if self._table_depth == 1:
                self._end_table_cell()
                self._end_table_row()
                self._parts.append("\n")
            self._table_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth > 0:
            return
        chunk = data.strip()
        if chunk:
            if self._table_depth > 0 and self._current_cell is not None:
                self._current_cell.append(chunk)
            elif self._table_depth == 0:
                self._parts.append(chunk)

    def get_text(self) -> str:
        raw = " ".join(self._parts)
        collapsed = re.sub(r"[ \t]+", " ", raw)
        collapsed = re.sub(r"\n\s*", "\n", collapsed)
        collapsed = re.sub(r" +\n", "\n", collapsed)
        return re.sub(r"\n{3,}", "\n\n", collapsed).strip()

    def _span_value(self, attrs: List, name: str) -> int:
        value = dict(attrs).get(name)
        try:
            return max(1, min(int(value or 1), 100))
        except ValueError:
            return 1

    def _start_table_row(self) -> None:
        self._end_table_cell()
        self._end_table_row()
        self._current_row = []
        self._new_rowspans = {}

    def _start_table_cell(self, attrs: List) -> None:
        if self._current_row is None:
            self._start_table_row()
        self._end_table_cell()
        self._current_cell = []
        self._current_colspan = self._span_value(attrs, "colspan")
        self._current_rowspan = self._span_value(attrs, "rowspan")

    def _end_table_cell(self) -> None:
        if self._current_cell is None or self._current_row is None:
            return

        col = len(self._current_row)
        while self._active_rowspans.get(col, 0) > 0:
            self._current_row.append("")
            col += 1

        text = re.sub(r"\s+", " ", " ".join(self._current_cell)).strip()
        for offset in range(self._current_colspan):
            self._current_row.append(text if offset == 0 else "")
            if self._current_rowspan > 1:
                self._new_rowspans[col + offset] = max(
                    self._new_rowspans.get(col + offset, 0),
                    self._current_rowspan - 1,
                )

        self._current_cell = None
        self._current_colspan = 1
        self._current_rowspan = 1

    def _end_table_row(self) -> None:
        if self._current_row is None:
            return

        self._end_table_cell()
        if self._active_rowspans:
            max_col = max(self._active_rowspans)
            while len(self._current_row) <= max_col:
                self._current_row.append("")

        if any(cell for cell in self._current_row):
            self._parts.append("| " + " | ".join(self._current_row) + " |")
            self._parts.append("\n")

        next_rowspans = {
            col: remaining - 1
            for col, remaining in self._active_rowspans.items()
            if remaining > 1
        }
        for col, remaining in self._new_rowspans.items():
            next_rowspans[col] = max(next_rowspans.get(col, 0), remaining)

        self._active_rowspans = next_rowspans
        self._new_rowspans = {}
        self._current_row = None


def _extract_html_text(html: str) -> str:
    extractor = _TextExtractor()
    extractor.feed(html)
    extractor.close()
    return extractor.get_text()


_MAX_FETCH_REDIRECTS = 5
_REDIRECT_STATUS_CODES = {301, 302, 303, 307, 308}


class _NoRedirectHandler(urllib.request.HTTPRedirectHandler):
    def redirect_request(self, req, fp, code, msg, headers, newurl):
        return None


_NO_REDIRECT_OPENER = urllib.request.build_opener(_NoRedirectHandler)


def _is_private_host(host: str) -> bool:
    try:
        ip = ipaddress.ip_address(host)
        return not ip.is_global
    except ValueError:
        pass

    try:
        infos = socket.getaddrinfo(host, None, type=socket.SOCK_STREAM)
    except Exception:
        return True
    if not infos:
        return True
    for info in infos:
        sockaddr = info[4]
        if not sockaddr:
            return True
        try:
            ip = ipaddress.ip_address(sockaddr[0])
        except ValueError:
            return True
        if not ip.is_global:
            return True
    return False


def _validate_public_fetch_url(raw_url: str) -> str:
    parsed = urlparse(raw_url)
    if parsed.scheme not in ("http", "https"):
        raise HTTPException(status_code=400, detail="only http/https allowed")
    if not parsed.hostname:
        raise HTTPException(status_code=400, detail="invalid url")
    if _is_private_host(parsed.hostname):
        raise HTTPException(status_code=400, detail="refusing to fetch private/internal host")
    return raw_url


def _open_public_url(raw_url: str, headers: Dict[str, str]):
    current_url = _validate_public_fetch_url(raw_url)
    for _ in range(_MAX_FETCH_REDIRECTS + 1):
        request = urllib.request.Request(current_url, headers=headers)
        try:
            return _NO_REDIRECT_OPENER.open(request, timeout=10)
        except urllib.error.HTTPError as e:
            if e.code not in _REDIRECT_STATUS_CODES:
                raise
            location = e.headers.get("Location")
            if not location:
                raise HTTPException(status_code=502, detail="redirect missing Location")
            current_url = _validate_public_fetch_url(urljoin(current_url, location))
    raise HTTPException(status_code=508, detail="too many redirects")


@app.post("/api/fetch-url")
async def fetch_url(req: FetchUrlRequest):
    def _do_fetch() -> Dict[str, str]:
        headers = {
            "User-Agent": "Mozilla/5.0 (compatible; ClaudeWeb/1.0)",
            "Accept": "text/html,application/xhtml+xml,text/plain;q=0.9,*/*;q=0.5",
        }
        with _open_public_url(req.url, headers) as resp:
            content_type = resp.headers.get("Content-Type", "") or ""
            charset = "utf-8"
            if "charset=" in content_type:
                charset = content_type.split("charset=", 1)[1].split(";")[0].strip()
            raw = resp.read(2 * 1024 * 1024)
        html = raw.decode(charset, errors="replace")
        title_match = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
        title = re.sub(r"\s+", " ", title_match.group(1)).strip() if title_match else req.url
        text = _extract_html_text(html)
        return {"title": title, "content": text}

    try:
        result = await asyncio.get_event_loop().run_in_executor(None, _do_fetch)
    except HTTPException:
        raise
    except urllib.error.HTTPError as e:
        raise HTTPException(status_code=502, detail=f"remote {e.code}")
    except urllib.error.URLError as e:
        raise HTTPException(status_code=502, detail=f"fetch failed: {e.reason}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    limit = max(500, min(req.max_chars or 10000, 50000))
    content = result["content"][:limit]
    return {
        "url": req.url,
        "title": result["title"] or req.url,
        "content": content,
        "truncated": len(result["content"]) > limit,
        "length": len(result["content"]),
    }


# ===== Enhanced Git Status Detail =====


def _parse_git_status_porcelain(lines: List[str]) -> List[dict]:
    """Parse git status --porcelain -z output into structured file list.

    In -z format, entries are separated by null bytes.
    Each entry: STATUS<SPACE>PATH1[<NULL>PATH2]
    STATUS is typically 2 chars (e.g. ' M', '??', 'MM') or 4 for renames (e.g. 'R100').
    PATH2 is only present for renames/copies.
    """
    result: List[dict] = []
    for line in lines:
        if not line:
            continue
        # Find the null byte first (in -z format it separates entries)
        null_idx = line.find('\0')
        if null_idx >= 0:
            entry = line[:null_idx]
            remaining = line[null_idx:]
        else:
            entry = line
            remaining = ''

        # Status is 4 chars for renames/copies (e.g. 'R100', 'C099'),
        # otherwise 2 chars (e.g. ' M', '??', 'MM').
        # After status there's always a space, then the path.
        if len(entry) >= 5 and entry[4] == ' ':
            # 4-char status (rename/copy with score)
            raw_status = entry[:4]
            rest = entry[5:]
            status = raw_status[0] + ' '
        elif len(entry) >= 3 and entry[2] == ' ':
            # 2-char status
            raw_status = entry[:2]
            rest = entry[3:]
            status = raw_status
        else:
            continue

        # For renames/copies, extract the second path from remaining
        renames = None
        if remaining.startswith('\0'):
            rename_path = remaining[1:].strip()
            if rename_path:
                renames = rename_path

        result.append({
            'status': status,
            'secondary': None,
            'filename': rest,
            'renames': renames,
        })
    return result


def _file_status_category(status: str) -> str:
    """Map git status code to category group."""
    first = status[0] if len(status) >= 1 else '?'
    second = status[1] if len(status) >= 2 else ' '
    if first == '?':
        return 'untracked'  # 未跟踪
    if first == 'D':
        return 'deleted'  # 已删除
    if first == 'R':
        return 'renamed'  # 已重命名
    if first == 'C':
        return 'copied'  # 已复制
    if first == 'A':
        return 'added'  # 已添加(新文件)
    if first == 'M' or first == 'T' or first == 'U':
        return 'modified'  # 已修改
    if first == ' ':
        if second == 'D':
            return 'deleted_staged'
        if second == 'M':
            return 'staged_modified'
        if second == 'A':
            return 'staged_added'
        if second == 'R':
            return 'staged_renamed'
        if second == 'C':
            return 'staged_copied'
        if second == '?':
            return 'staged_untracked'
        return 'staged_modified'
    return 'other'


_STATUS_GROUPS = {
    'staged_modified': {'label': '已暂存的修改', 'icon': '📝', 'color': '#3b82f6', 'key': 'staged_modified'},
    'staged_added': {'label': '已暂存的新文件', 'icon': '📄', 'color': '#3b82f6', 'key': 'staged_added'},
    'staged_renamed': {'label': '已暂存的 renaming', 'icon': '🏷️', 'color': '#3b82f6', 'key': 'staged_renamed'},
    'staged_copied': {'label': '已暂存的复制', 'icon': '📋', 'color': '#3b82f6', 'key': 'staged_copied'},
    'staged_untracked': {'label': '已暂存的未跟踪文件', 'icon': '📌', 'color': '#6366f1', 'key': 'staged_untracked'},
    'modified': {'label': '已修改', 'icon': '✏️', 'color': '#f59e0b', 'key': 'modified'},
    'deleted': {'label': '已删除', 'icon': '🗑️', 'color': '#ef4444', 'key': 'deleted'},
    'deleted_staged': {'label': '已暂存的删除', 'icon': '🗑️', 'color': '#ef4444', 'key': 'deleted_staged'},
    'added': {'label': '新文件', 'icon': '✨', 'color': '#10b981', 'key': 'added'},
    'renamed': {'label': '已重命名', 'icon': '🏷️', 'color': '#8b5cf6', 'key': 'renamed'},
    'copied': {'label': '已复制', 'icon': '📋', 'color': '#06b6d4', 'key': 'copied'},
    'untracked': {'label': '未跟踪', 'icon': '❓', 'color': '#6b7280', 'key': 'untracked'},
    'other': {'label': '其他', 'icon': '⚠️', 'color': '#6b7280', 'key': 'other'},
}


@app.get("/api/git/status-detail")
async def git_status_detail(cwd: str = Query(default="")):
    """Get detailed git status grouped by change type, with staged files info."""
    git_cwd = os.path.expanduser(cwd) if cwd else "."
    if not os.path.isdir(git_cwd):
        return {"branch": "", "dirty": 0, "available": False}

    # Get branch
    branch_proc = await asyncio.create_subprocess_exec(
        "git", "-C", git_cwd, "rev-parse", "--abbrev-ref", "HEAD",
        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    try:
        branch_out, _ = await asyncio.wait_for(branch_proc.communicate(), timeout=5)
    except asyncio.TimeoutError:
        branch_proc.kill()
        return {"branch": "", "dirty": 0, "available": False}
    branch = branch_out.decode("utf-8", errors="replace").strip() if branch_proc.returncode == 0 else ""

    # Get staged files
    staged_proc = await asyncio.create_subprocess_exec(
        "git", "-C", git_cwd, "diff", "--cached", "--name-only",
        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    try:
        staged_out, _ = await asyncio.wait_for(staged_proc.communicate(), timeout=5)
    except asyncio.TimeoutError:
        staged_proc.kill()
        return {"branch": branch, "dirty": 0, "available": True, "files": [], "staged": []}
    staged_files = []
    if staged_proc.returncode == 0 and staged_out.strip():
        staged_files = [f.strip() for f in staged_out.decode("utf-8", errors="replace").splitlines() if f.strip()]

    # Get detailed porcelain status
    porcelain_proc = await asyncio.create_subprocess_exec(
        "git", "-C", git_cwd, "status", "--porcelain", "-z",
        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    try:
        porcelain_out, _ = await asyncio.wait_for(porcelain_proc.communicate(), timeout=5)
    except asyncio.TimeoutError:
        porcelain_proc.kill()
        return {"branch": branch, "dirty": 0, "available": True, "files": [], "staged": staged_files}
    if porcelain_proc.returncode != 0:
        return {"branch": branch, "dirty": 0, "available": True, "files": [], "staged": staged_files}

    raw = porcelain_out.decode("utf-8", errors="replace")
    # Split by null byte
    entries = raw.split('\0') if '\0' in raw else raw.splitlines()
    lines = [e for e in entries if e.strip()]
    files = _parse_git_status_porcelain(lines)

    # Group by category
    groups: Dict[str, List[dict]] = {}
    for f in files:
        cat = _file_status_category(f['status'])
        groups.setdefault(cat, []).append(f)

    # Build ordered result
    ordered_files = []
    for cat_key, info in _STATUS_GROUPS.items():
        group_files = groups.get(cat_key, [])
        if not group_files:
            continue
        ordered_files.append({
            'category': info['label'],
            'icon': info['icon'],
            'color': info['color'],
            'key': cat_key,
            'files': group_files,
        })

    # Count dirty
    dirty = len(files)

    # Get remote info
    remote_proc = await asyncio.create_subprocess_exec(
        "git", "-C", git_cwd, "remote", "get-url", "origin",
        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    try:
        remote_out, _ = await asyncio.wait_for(remote_proc.communicate(), timeout=3)
    except asyncio.TimeoutError:
        remote_out = b""
    remote_url = remote_out.decode("utf-8", errors="replace").strip() if remote_proc.returncode == 0 else ""

    return {
        "branch": branch,
        "dirty": dirty,
        "available": True,
        "files": ordered_files,
        "staged": staged_files,
        "remote_url": remote_url,
    }


# ===== File Explorer & Git Operations =====


class DirPickerRequest(BaseModel):
    cwd: str


class FileContentRequest(BaseModel):
    path: str
    max_lines: Optional[int] = 10000


class FileSaveRequest(BaseModel):
    path: str
    content: str


class GitRunRequest(BaseModel):
    cwd: str
    command: str
    args: Optional[List[str]] = None


class GitDiffRequest(BaseModel):
    path: str
    cwd: str
    cached: Optional[bool] = False


class GitLogRequest(BaseModel):
    cwd: str
    limit: Optional[int] = 50


# Whitelisted git command patterns: (main_cmd, allowed_subcmds)
_GIT_CMD_WHITELIST = {
    "init": [],
    "clone": [],
    "status": [],
    "add": [],
    "commit": [],
    "push": [],
    "pull": [],
    "fetch": [],
    "branch": [],
    "checkout": [],
    "switch": [],
    "merge": [],
    "rebase": [],
    "log": [],
    "diff": [],
    "stash": [],
    "reset": [],
    "revert": [],
    "tag": [],
    "remote": [],
    "rm": [],
    "mv": [],
}


def _sanitize_path(p: str) -> Path:
    resolved = Path(os.path.expanduser(p)).resolve()
    return resolved


def _validate_path_in_dir(target: Path, parent: Path) -> bool:
    try:
        target.relative_to(parent)
        return True
    except ValueError:
        return False


def _safe_git_run(cwd: str, *args: str) -> Optional[dict]:
    """Run a git command safely. Returns dict with stdout/stderr/rc or None on error."""
    try:
        proc = asyncio.get_event_loop().run_in_executor(
            None,
            _git_run_sync,
            os.path.expanduser(cwd),
            args,
        )
        return asyncio.get_event_loop().run_until_complete(proc)
    except Exception:
        return None


def _git_run_sync(cwd: str, args: tuple) -> dict:
    try:
        proc = subprocess.run(
            ["git", "-C", cwd] + list(args),
            capture_output=True,
            timeout=15,
        )
        return {
            "stdout": proc.stdout.decode("utf-8", errors="replace"),
            "stderr": proc.stderr.decode("utf-8", errors="replace"),
            "returncode": proc.returncode,
        }
    except subprocess.TimeoutExpired:
        return {"stdout": "", "stderr": "timeout", "returncode": -1}
    except FileNotFoundError:
        return {"stdout": "", "stderr": "git not found", "returncode": -1}
    except Exception as e:
        return {"stdout": "", "stderr": str(e), "returncode": -1}


@app.post("/api/dir-picker")
async def dir_picker(req: DirPickerRequest):
    """Return subdirectories of the given cwd for the folder picker."""
    base = _sanitize_path(req.cwd)
    if not base.is_dir():
        return {"dirs": []}
    dirs = []
    try:
        for entry in sorted(base.iterdir()):
            if entry.is_dir() and entry.name not in IGNORED_DIRS and not entry.name.startswith("."):
                dirs.append({"name": entry.name, "path": str(entry)})
    except OSError:
        pass
    return {"cwd": str(base), "dirs": dirs}


@app.post("/api/tree")
async def get_tree(req: FileContentRequest):
    """Get immediate children of a directory (lazy-loading tree)."""
    base = _sanitize_path(req.path)
    if not base.is_dir():
        return {"path": str(base), "children": []}
    children = []
    try:
        for entry in sorted(base.iterdir(), key=lambda e: (not e.is_dir(), e.name.lower())):
            if entry.name.startswith(".") and entry.name not in (".github", ".vscode", ".idea"):
                continue
            if entry.name in IGNORED_DIRS:
                continue
            info: dict = {"name": entry.name, "type": "dir" if entry.is_dir() else "file"}
            if entry.is_file():
                try:
                    stat = entry.stat()
                    info["size"] = stat.st_size
                    info["modified"] = stat.st_mtime
                except OSError:
                    pass
            children.append(info)
    except OSError:
        pass
    return {"path": str(base), "children": children}


@app.get("/api/file-content")
async def get_file_content(path: str = Query(...), max_lines: int = Query(default=10000)):
    """Read a file's content with line numbers."""
    target = _sanitize_path(path)
    if not target.is_file():
        raise HTTPException(status_code=404, detail="file not found")
    try:
        raw = target.read_bytes()
    except OSError as e:
        raise HTTPException(status_code=400, detail=str(e))
    if len(raw) > 5 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="file too large (max 5MB)")
    text = _decode_text_upload(raw) if b"\x00" in raw[:8192] else raw.decode("utf-8", errors="replace")
    all_lines = text.split("\n")
    total = len(all_lines)
    display_lines = all_lines[:max_lines]
    if len(all_lines) > max_lines:
        display_lines.append(f"... ({total - max_lines} more lines truncated ...)")
    lang = _detect_language(str(target))
    return {
        "path": str(target),
        "content": text,
        "lines": [{"num": i + 1, "text": l} for i, l in enumerate(display_lines)],
        "lines_total": total,
        "language": lang,
        "size": len(raw),
    }


def _detect_language(filename: str) -> str:
    ext = Path(filename).suffix.lower().lstrip(".")
    EXT_LANG_MAP = {
        "py": "python", "js": "javascript", "ts": "typescript", "tsx": "typescript", "jsx": "javascript",
        "go": "go", "rs": "rust", "rb": "ruby", "java": "java", "c": "c", "cpp": "cpp", "h": "c",
        "cs": "csharp", "php": "php", "swift": "swift", "kt": "kotlin", "scala": "scala",
        "sh": "bash", "bash": "bash", "zsh": "bash", "ps1": "powershell", "psm1": "powershell",
        "html": "html", "css": "css", "scss": "scss", "less": "less",
        "json": "json", "yaml": "yaml", "yml": "yaml", "toml": "toml", "xml": "xml",
        "md": "markdown", "sql": "sql", "dockerfile": "dockerfile",
        "lua": "lua", "r": "r", "m": "matlab", "pl": "perl",
        "proto": "protobuf", "graphql": "graphql",
    }
    return EXT_LANG_MAP.get(ext, "")


@app.post("/api/file-save")
async def save_file(req: FileSaveRequest):
    """Save content to a file."""
    target = _sanitize_path(req.path)
    if not target.parent.is_dir():
        raise HTTPException(status_code=400, detail="parent directory not found")
    if len(req.content) > 10 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="content too large")
    try:
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(req.content, encoding="utf-8")
    except OSError as e:
        raise HTTPException(status_code=400, detail=str(e))
    return {"ok": True, "path": str(target)}


@app.post("/api/git/discard")
async def git_discard(req: GitRunRequest):
    """Discard changes for selected files (checkout for tracked, clean for untracked)."""
    cwd = os.path.expanduser(req.cwd)
    if not os.path.isdir(cwd):
        raise HTTPException(status_code=400, detail="cwd not found")
    filename = req.command.strip()  # the filename to discard
    # Sanitize path
    try:
        target = Path(filename).resolve()
        target.relative_to(Path(cwd).resolve())
    except (ValueError, OSError):
        raise HTTPException(status_code=400, detail="invalid filename")
    # First try git checkout for tracked files
    proc = await asyncio.create_subprocess_exec(
        "git", "-C", cwd, "checkout", "--", filename,
        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    _, _ = await proc.communicate()
    if proc.returncode == 0:
        return {"ok": True, "stdout": "", "stderr": ""}
    # If checkout failed (untracked), try git clean
    proc = await asyncio.create_subprocess_exec(
        "git", "-C", cwd, "clean", "-f", "--", filename,
        stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
    )
    _, _ = await proc.communicate()
    if proc.returncode == 0:
        return {"ok": True, "stdout": "", "stderr": ""}
    # Last resort: rm -f via subprocess
    try:
        rm_proc = await asyncio.create_subprocess_exec(
            "rm", "-f", filename,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
        )
        _, _ = await rm_proc.communicate()
        if rm_proc.returncode == 0:
            return {"ok": True, "stdout": "", "stderr": ""}
    except Exception:
        pass
    raise HTTPException(status_code=400, detail="discard failed: file may not exist or is not cleanable")


@app.post("/api/git/run")
async def git_run(req: GitRunRequest):
    """Execute a git command from the whitelist."""
    cwd = os.path.expanduser(req.cwd)
    if not os.path.isdir(cwd):
        raise HTTPException(status_code=400, detail="cwd not found")
    parts = req.command.strip().split()
    if not parts:
        raise HTTPException(status_code=400, detail="empty command")
    main_cmd = parts[0].lower()
    if main_cmd != "git":
        raise HTTPException(status_code=403, detail=f"command not allowed: {main_cmd}")
    sub_parts = parts[1:]
    if not sub_parts:
        raise HTTPException(status_code=400, detail="missing subcommand")
    subcmd = sub_parts[0].lower()
    if subcmd not in _GIT_CMD_WHITELIST:
        raise HTTPException(status_code=403, detail=f"git subcommand not allowed: {subcmd}")
    # Block dangerous operations on remote
    if subcmd in ("push",) and "--force" in " ".join(parts[2:]):
        raise HTTPException(status_code=403, detail="force push is not allowed via UI")
    git_args = parts[1:]
    result = await asyncio.to_thread(_git_run_sync, cwd, tuple(git_args))
    return result


@app.get("/api/git/diff")
async def git_diff(path: str = Query(...), cwd: str = Query(default=""), cached: bool = Query(default=False)):
    """Get git diff for a specific file."""
    git_cwd = os.path.expanduser(cwd) if cwd else "."
    if not os.path.isdir(git_cwd):
        return {"file": "", "diff_lines": []}
    try:
        target = _sanitize_path(path)
        rel_path = str(target.relative_to(_sanitize_path(git_cwd))) if _sanitize_path(git_cwd) in target.parents or _sanitize_path(git_cwd) == target else target.name
    except (ValueError, OSError):
        rel_path = Path(path).name
    if cached:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", git_cwd, "diff", "--cached", "-U0", "--", rel_path,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
        )
    else:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", git_cwd, "diff", "-U0", "--", rel_path,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
        )
    try:
        stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=10)
    except asyncio.TimeoutError:
        proc.kill()
        await proc.wait()
        return {"file": rel_path, "diff_lines": []}
    if proc.returncode != 0:
        return {"file": rel_path, "diff_lines": []}
    output = stdout.decode("utf-8", errors="replace").strip()
    diff_lines = _parse_git_diff_lines(output)
    return {"file": rel_path, "diff_lines": diff_lines}


def _parse_git_diff_lines(output: str) -> List[dict]:
    """Parse unified diff output into line-by-line entries with line numbers."""
    if not output:
        return []
    lines = output.split("\n")
    result: List[dict] = []
    old_line = 0
    new_line = 0
    for line in lines:
        if line.startswith("diff ") or line.startswith("index ") or line.startswith("--- ") or line.startswith("+++ "):
            continue
        if line.startswith("@@ "):
            # Parse @@ -old_start,old_count +new_start,new_count @@
            import re
            m = re.search(r"-(\d+)(?:,\d+)?\+(\d+)(?:,\d+)?", line)
            if m:
                old_line = int(m.group(1))
                new_line = int(m.group(2))
            continue
        if not line:
            continue
        if line[0] == "+":
            result.append({"line_old": None, "line_new": new_line, "content": "+ " + line[1:], "type": "add"})
            new_line += 1
        elif line[0] == "-":
            result.append({"line_old": old_line, "line_new": None, "content": "- " + line[1:], "type": "remove"})
            old_line += 1
        elif line[0] == " ":
            result.append({"line_old": old_line, "line_new": new_line, "content": "  " + line[1:], "type": "context"})
            old_line += 1
            new_line += 1
        else:
            # Binary or other
            result.append({"line_old": None, "line_new": None, "content": line, "type": "other"})
    return result


@app.get("/api/git/log")
async def git_log(cwd: str = Query(default=""), limit: int = Query(default=50)):
    """Get recent git log with graph."""
    git_cwd = os.path.expanduser(cwd) if cwd else "."
    if not os.path.isdir(git_cwd):
        return {"commits": [], "graph": ""}
    try:
        proc = await asyncio.create_subprocess_exec(
            "git", "-C", git_cwd, "log", "--graph", f"--max-count={limit}",
            "--oneline", "--format=%h %s %cd", "--date=relative",
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=10)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return {"commits": [], "graph": ""}
        if proc.returncode != 0:
            return {"commits": [], "graph": ""}
    except Exception:
        return {"commits": [], "graph": ""}
    raw = stdout.decode("utf-8", errors="replace")
    commits = []
    graph_lines = []
    for line in raw.splitlines():
        # Parse: graph_hash hash message relative_date
        # Example: * a1b2c3d feat: add new feature (2 hours ago)
        stripped = line.strip()
        if not stripped:
            continue
        graph_lines.append(line)
        # Try to extract hash and message
        # Remove leading graph chars
        cleaned = re.sub(r'^[\s\|\-\+\*\/\n]+', '', stripped)
        parts = cleaned.split(" ", 1)
        if len(parts) >= 2:
            hash_val = parts[0]
            message = parts[1]
            # Try to extract time from end
            time_match = re.search(r'\(([^)]+)\)$', message)
            time_str = ""
            msg_clean = message
            if time_match:
                time_str = time_match.group(1)
                msg_clean = message[:time_match.start()].rstrip()
            commits.append({"hash": hash_val, "message": msg_clean, "time": time_str, "graph": stripped})
        elif len(parts) == 1:
            graph = parts[0]
            if commits:
                commits[-1]["graph"] = (commits[-1].get("graph", "") + "\n" + graph).strip()
    return {"commits": commits, "graph": "\n".join(graph_lines[:10])}


# ===== CWD History =====

@app.get("/api/cwd-history")
async def list_cwd_history():
    with db_connect() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS cwd_history (
                id TEXT PRIMARY KEY,
                path TEXT NOT NULL UNIQUE,
                created_at REAL NOT NULL
            )
            """
        )
        rows = conn.execute(
            "SELECT path, created_at FROM cwd_history ORDER BY created_at DESC LIMIT 50"
        ).fetchall()
    return [{"path": r["path"], "created_at": r["created_at"]} for r in rows]


@app.post("/api/cwd-history")
async def upsert_cwd_history(path: str = Query(...)):
    path = path.strip()
    if not path:
        return {"ok": True}
    now = time.time()
    with db_connect() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS cwd_history (
                id TEXT PRIMARY KEY,
                path TEXT NOT NULL UNIQUE,
                created_at REAL NOT NULL
            )
            """
        )
        existing = conn.execute(
            "SELECT id FROM cwd_history WHERE path = ?", (path,)
        ).fetchone()
        if existing:
            conn.execute(
                "UPDATE cwd_history SET created_at = ? WHERE path = ?", (now, path)
            )
        else:
            conn.execute(
                "INSERT INTO cwd_history (id, path, created_at) VALUES (?, ?, ?)",
                (uuid.uuid4().hex, path, now),
            )
    return {"ok": True}


@app.delete("/api/cwd-history/{path:path}")
async def delete_cwd_history(path: str):
    decoded = urllib.request.unquote(path)
    with db_connect() as conn:
        conn.execute("DELETE FROM cwd_history WHERE path = ?", (decoded,))
    return {"ok": True}


@app.post("/api/cwd-history/clear")
async def clear_cwd_history():
    with db_connect() as conn:
        conn.execute("DELETE FROM cwd_history")
    return {"ok": True}


# ===== Config (settings / hooks / permissions / skills) =====

_CLAUDE_USER_DIR = Path.home() / ".claude"
_CLAUDE_SETTINGS_FILE = _CLAUDE_USER_DIR / "settings.json"
_SKILLS_DIR = _CLAUDE_USER_DIR / "skills"


def _read_json_file(path: Path, default: Optional[dict] = None) -> dict:
    if not path.exists():
        return default or {}
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return default or {}


def _write_json_file(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def _resolve_config_scope(scope: str, cwd: Optional[str]) -> Path:
    normalized = (scope or "local").strip().lower()
    if normalized == "user":
        return _CLAUDE_SETTINGS_FILE
    # "local" or "project" → .claude/settings.json in project
    project_dir = Path(os.path.expanduser(cwd)).resolve() if cwd else Path.cwd()
    return project_dir / ".claude" / "settings.json"


def _unflatten_hooks(settings: dict) -> List[dict]:
    hooks_config = settings.get("hooks", {}) or {}
    rows: List[dict] = []
    for event_name, event_groups in (hooks_config or {}).items():
        if not isinstance(event_groups, list):
            continue
        for group in event_groups:
            if not isinstance(group, dict):
                continue
            matcher = group.get("matcher", "")
            inner_hooks = group.get("hooks", []) or []
            for h in inner_hooks:
                if not isinstance(h, dict):
                    continue
                rows.append({
                    "event": event_name,
                    "matcher": matcher,
                    "type": h.get("type", "command"),
                    "command": h.get("command", ""),
                    "enabled": h.get("enabled", True),
                })
    return rows


def _flatten_hooks(rows: List[dict]) -> dict:
    grouped: Dict[str, List[dict]] = {}
    for row in rows:
        event = row.get("event", "PreToolUse")
        hook = {
            "matcher": row.get("matcher", ""),
            "type": row.get("type", "command"),
            "command": row.get("command", ""),
            "enabled": row.get("enabled", True),
        }
        grouped.setdefault(event, []).append(hook)
    return {"hooks": grouped}


def _flatten_permissions(settings: dict) -> dict:
    return settings.get("permissions", {}) or {}


# Skill definitions for SKILL.md frontmatter parsing
_SKILL_FM_RE = re.compile(r"^---\n(.*?)\n---", re.DOTALL)


def _parse_skill_frontmatter(md: str) -> tuple[Optional[dict], str]:
    m = _SKILL_FM_RE.match(md)
    fm = {}
    body = md
    if m:
        try:
            fm = json.loads(m.group(1))
        except json.JSONDecodeError:
            pass
        body = md[m.end():].lstrip("\n")
    return fm, body


def _scan_skills_dir(skills_dir: Path) -> List[dict]:
    if not skills_dir.is_dir():
        return []
    items: List[dict] = []
    for entry in sorted(skills_dir.iterdir()):
        if not entry.is_dir():
            continue
        skill_md = entry / "SKILL.md"
        if not skill_md.exists():
            continue
        name = entry.name
        try:
            raw = skill_md.read_text(encoding="utf-8", errors="replace")
        except OSError:
            continue
        fm, body = _parse_skill_frontmatter(raw)
        # Check disabled via SKILL.md.disabled pattern: if parent dir has no SKILL.md
        # but was scanned — disabled skills are .disabled dirs or have SKILL.md.disabled
        desc = (fm.get("description", "") or body.split("\n")[0][:200]).strip()
        enabled = True
        items.append({
            "name": name,
            "description": desc,
            "enabled": enabled,
            "marketplace": "@local",
            "error": None,
        })
    return items


class ConfigSettingsRequest(BaseModel):
    scope: str = "local"
    settings: dict


class SkillToggleRequest(BaseModel):
    enabled: bool


class SkillTranslateItem(BaseModel):
    name: str
    description: str


class SkillTranslateRequest(BaseModel):
    items: List[SkillTranslateItem]


@app.get("/api/config/settings")
async def get_config_settings(scope: str = Query(default="local"), cwd: Optional[str] = Query(default=None)):
    target = _resolve_config_scope(scope, cwd)
    settings = _read_json_file(target, {})
    return {
        "path": str(target),
        "exists": target.exists(),
        "settings": settings,
        "scope": scope,
    }


@app.post("/api/config/settings")
async def save_config_settings(req: ConfigSettingsRequest):
    target = _resolve_config_scope(req.scope, None)
    existing = _read_json_file(target, {})
    # Merge: update hooks and permissions only
    if "hooks" in req.settings:
        existing["hooks"] = req.settings["hooks"]
    if "permissions" in req.settings:
        existing["permissions"] = req.settings["permissions"]
    _write_json_file(target, existing)
    return {
        "settings": existing,
        "path": str(target),
        "exists": True,
    }


@app.get("/api/config/skills")
async def list_config_skills(cwd: Optional[str] = Query(default=None)):
    skills_dir = _SKILLS_DIR
    # Also check project-level .claude/skills
    project_skills = None
    if cwd:
        project_dir = Path(os.path.expanduser(cwd)).resolve()
        project_skills = project_dir / ".claude" / "skills"

    items = _scan_skills_dir(skills_dir)
    skills_dir_str = str(_SKILLS_DIR)

    if project_skills and project_skills.is_dir():
        project_items = _scan_skills_dir(project_skills)
        # Prefix project skills with project name
        proj_name = project_dir.name
        for it in project_items:
            it["name"] = f"{proj_name}/{it['name']}"
            it["marketplace"] = f"@project({proj_name})"
        items.extend(project_items)
        skills_dir_str = f"{skills_dir_str} + {project_skills}"

    return {
        "skills": items,
        "skills_dir": skills_dir_str,
    }


def _resolve_skill_path(skill_name: str) -> Path | None:
    """Find the directory containing a skill's SKILL.md by name.
    Returns the skill directory path, or None if not found."""
    # Strip project prefix if present (e.g. "myproject/my-skill" → "my-skill")
    parts = skill_name.split("/", 1)
    short = parts[-1]

    # Check user skills dir first
    user_skill = _SKILLS_DIR / short
    if user_skill.is_dir():
        return user_skill

    # Check project-level .claude/skills dirs (scan from cwd or common projects)
    for candidate in _SKILLS_DIR.parent.glob("*"):
        if not candidate.is_dir():
            continue
        proj_skill = candidate / ".claude" / "skills" / short
        if proj_skill.is_dir():
            return proj_skill

    # Fallback: return user dir path anyway (for error messages)
    return _SKILLS_DIR / short


def _toggle_skill_file(skill_dir: Path, enabled: bool) -> bool:
    """Rename SKILL.md ↔ SKILL.md.disabled. Returns True if a file was changed."""
    skill_file = skill_dir / "SKILL.md"
    disabled_file = skill_dir / "SKILL.md.disabled"
    if enabled and disabled_file.exists():
        disabled_file.rename(skill_file)
        return True
    elif not enabled and skill_file.exists():
        skill_file.rename(disabled_file)
        return True
    return False


@app.patch("/api/config/skills/{skill_name}")
async def toggle_skill(skill_name: str, req: SkillToggleRequest):
    """Toggle a skill by renaming SKILL.md ↔ SKILL.md.disabled in its directory."""
    skill_dir = _resolve_skill_path(skill_name)
    if skill_dir is None or not skill_dir.is_dir():
        raise HTTPException(status_code=404, detail=f"skill '{skill_name}' not found")
    _toggle_skill_file(skill_dir, req.enabled)
    return {"ok": True}


@app.get("/api/config/skills/{skill_name}/source")
async def get_skill_source(skill_name: str):
    skill_dir = _resolve_skill_path(skill_name)
    if skill_dir is None or not skill_dir.is_dir():
        raise HTTPException(status_code=404, detail=f"skill file not found for {skill_name}")
    skill_file = skill_dir / "SKILL.md"
    disabled_file = skill_dir / "SKILL.md.disabled"

    target = skill_file if skill_file.exists() else disabled_file
    if not target.exists():
        raise HTTPException(status_code=404, detail=f"no SKILL.md found for skill '{skill_name}'")

    content = target.read_text(encoding="utf-8", errors="replace")
    return {"name": skill_name, "content": content}


@app.post("/api/config/skills/translate")
async def translate_skills(req: SkillTranslateRequest):
    """Translate skill descriptions to Chinese via Claude CLI."""
    if not req.items:
        return {"translations": {}}

    # Build prompt for each skill
    items_list = [f"- {item.name}: {item.description}" for item in req.items]
    prompt = (
        "请将以下 skill 描述翻译为中文（每个一行，格式为 \"name: 翻译结果\"），"
        "不要添加额外内容。如果没有描述或不是英文则原样返回原文。\n\n"
        + "\n".join(items_list)
    )

    try:
        proc = await asyncio.create_subprocess_exec(
            *claude_cli_argv("-p", prompt, "--output-format", "text", "--model", "haiku"),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            return {"translations": {}}
    except Exception:
        return {"translations": {}}

    text = stdout.decode("utf-8", errors="replace").strip()
    translations: Dict[str, str] = {}
    for line in text.splitlines():
        line = line.strip()
        if not line or ": " not in line:
            continue
        name_part, _, desc_part = line.partition(": ")
        name_part = name_part.strip()
        # Match against the original item names
        for item in req.items:
            skill_name = item.name.split("/", 1)[-1] if "/" in item.name else item.name
            if name_part == item.name or name_part == skill_name:
                translations[item.name] = desc_part.strip()
                break

    return {"translations": translations}


@app.get("/api/version")
async def get_version():
    return {"version": __version__}


@app.get("/changelog.json")
async def get_changelog():
    path = STATIC_DIR / "changelog.json"
    if not path.exists():
        return {"releases": []}
    return FileResponse(path, media_type="application/json")


@app.get("/")
async def index():
    return FileResponse(STATIC_DIR / "index.html")


def _check_claude_cli() -> Optional[str]:
    """Return claude CLI version string if available, else None."""
    import subprocess

    command = resolve_claude_cli_command()
    if command is None:
        return None
    try:
        result = subprocess.run(
            claude_cli_argv("--version", allow_batch_shim=True),
            capture_output=True,
            text=True,
            timeout=5,
        )
        if result.returncode == 0:
            return result.stdout.strip() or result.stderr.strip() or "unknown"
    except (subprocess.TimeoutExpired, OSError, ClaudeCliResolutionError):
        pass
    return "unknown"


def main():
    """CLI entry point for `claude-web` command."""
    import argparse
    import sys
    import uvicorn

    parser = argparse.ArgumentParser(description="Claude Code Web - Web UI for Claude Code CLI")
    parser.add_argument("--port", "-p", type=int, default=int(os.environ.get("PORT", "8765")), help="Port to listen on (default: 8765)")
    parser.add_argument("--host", type=str, default="127.0.0.1", help="Host to bind (default: 127.0.0.1)")
    parser.add_argument("--open", action="store_true", help="Open browser after starting")
    parser.add_argument("--version", "-v", action="store_true", help="Show version")
    parser.add_argument("--skip-cli-check", action="store_true", help="Skip claude CLI availability check on startup")
    args = parser.parse_args()

    if args.version:
        print(f"claude-web {__version__}")
        return

    print(f"Claude Code Web v{__version__}")
    print(f"  → http://{args.host}:{args.port}")
    print(f"  → Data: {_DATA_DIR}")

    if not args.skip_cli_check:
        claude_version = _check_claude_cli()
        if claude_version is None:
            print()
            print("  ✗ claude CLI not found in PATH", file=sys.stderr)
            print("    claude-web wraps the Claude Code CLI — install it first:", file=sys.stderr)
            print("      npm install -g @anthropic-ai/claude-code", file=sys.stderr)
            print("    Then run `claude` once to log in. Docs: https://docs.claude.com/claude-code", file=sys.stderr)
            print("    (Use --skip-cli-check to bypass this check.)", file=sys.stderr)
            print()
            sys.exit(1)
        print(f"  → Claude CLI: {claude_version}")

    _LOCAL_HOSTS = {"127.0.0.1", "localhost", "::1"}
    if args.host not in _LOCAL_HOSTS:
        print()
        print(f"  ⚠️  WARNING: binding to {args.host} exposes the server beyond localhost.", file=sys.stderr)
        print("     claude-web has NO built-in authentication. Anyone who can reach this", file=sys.stderr)
        print("     address can run commands, read your files, and burn your Claude quota.", file=sys.stderr)
        print("     Only use --host on a trusted network (e.g. tailscale, VPN, SSH tunnel).", file=sys.stderr)

    print()

    if args.open:
        import webbrowser
        import threading
        threading.Timer(1.5, lambda: webbrowser.open(f"http://{args.host}:{args.port}")).start()

    uvicorn.run(app, host=args.host, port=args.port)


if __name__ == "__main__":
    main()
